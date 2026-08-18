"""Agent 运行器：以子进程无人值守地驱动 Claude Code / Codex CLI 在 ppt-master 仓库内执行工作流。

三种运行器：
- ClaudeRunner：`claude -p <prompt> --output-format stream-json --verbose --dangerously-skip-permissions`
- CodexRunner ：`codex exec --json --dangerously-bypass-approvals-and-sandbox --skip-git-repo-check -C <repo> <prompt>`
- MockRunner  ：不调用任何 Agent，用 python-pptx 直接产出占位 PPTX（链路联调 / 单测 / 无 Agent 环境）

统一契约：`run()` 逐行回调 stdout（供进度解析与日志落盘），支持超时与取消（杀进程树），
返回 RunResult（退出码 / 最终文本 / 用量）。
"""
from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import threading
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable

from app.core.config import get_settings
from app.core.logging import get_logger

logger = get_logger(__name__)

_IS_WIN = os.name == "nt"


# ---------------------------------------------------------------- 探测 ----
@dataclass
class AgentInfo:
    key: str
    label: str
    available: bool
    bin: str | None = None
    note: str = ""


def _candidates(name: str, configured: str) -> list[str]:
    cands: list[str] = []
    if configured:
        cands.append(configured)
    w = shutil.which(name)
    if w:
        cands.append(w)
    home = Path.home()
    cands += [str(home / ".local" / "bin" / f"{name}.exe"), str(home / ".local" / "bin" / name)]
    if _IS_WIN:
        appdata = os.environ.get("APPDATA", "")
        if appdata:
            cands += [os.path.join(appdata, "npm", f"{name}.cmd")]
    seen, out = set(), []
    for c in cands:
        if c and c not in seen and os.path.exists(c):
            seen.add(c)
            out.append(c)
    return out


def _version_ok(binary: str, expect: str, timeout: float = 20.0) -> tuple[bool, str]:
    try:
        r = subprocess.run([binary, "--version"], capture_output=True, text=True, timeout=timeout,
                           encoding="utf-8", errors="replace")
        out = (r.stdout or "") + (r.stderr or "")
        return (r.returncode == 0 and expect.lower() in out.lower()), out.strip().splitlines()[0] if out.strip() else ""
    except Exception as e:  # noqa: BLE001
        return False, str(e)


_detect_cache: dict[str, tuple[float, list[AgentInfo]]] = {}


def detect_agents(force: bool = False) -> list[AgentInfo]:
    """探测本机可用 Agent CLI（结果缓存 60s；PATH 上的失效垫片会被 --version 探测过滤掉）。"""
    now = time.time()
    if not force and "agents" in _detect_cache and now - _detect_cache["agents"][0] < 60:
        return _detect_cache["agents"][1]
    s = get_settings()
    infos: list[AgentInfo] = []

    claude = AgentInfo("claude", "Claude Code CLI", False, note="未找到可用的 claude 可执行文件")
    for c in _candidates("claude", s.pptmaster_claude_bin):
        ok, ver = _version_ok(c, "claude")
        if ok:
            claude = AgentInfo("claude", "Claude Code CLI", True, c, ver)
            break
    infos.append(claude)

    codex = AgentInfo("codex", "Codex CLI", False, note="未找到可用的 codex 可执行文件")
    for c in _candidates("codex", s.pptmaster_codex_bin):
        ok, ver = _version_ok(c, "codex")
        if ok:
            codex = AgentInfo("codex", "Codex CLI", True, c, ver)
            break
    infos.append(codex)

    infos.append(AgentInfo("mock", "Mock（不调用 Agent，仅生成占位 PPTX 验证链路）", True, None, "内置"))
    _detect_cache["agents"] = (now, infos)
    return infos


def select_agent(requested: str, available: list[AgentInfo], default_agent: str = "auto") -> AgentInfo:
    """从给定能力快照解析 auto|claude|codex|mock；不可用则抛 ValueError。"""
    agents = {a.key: a for a in available}
    want = requested or "auto"
    if want == "auto":
        want = default_agent or "auto"
    if want == "auto":
        for k in ("claude", "codex", "mock"):
            if k in agents and agents[k].available:
                return agents[k]
    info = agents.get(want)
    if info is None:
        raise ValueError(f"未知 Agent：{requested}")
    if not info.available:
        raise ValueError(f"Agent {want} 不可用：{info.note}")
    return info


def resolve_agent(requested: str) -> AgentInfo:
    """在当前进程内探测并解析 Agent。"""
    s = get_settings()
    return select_agent(requested, detect_agents(), s.pptmaster_default_agent)


# ---------------------------------------------------------------- 运行 ----
@dataclass
class RunResult:
    returncode: int
    timed_out: bool = False
    canceled: bool = False
    final_text: str = ""            # Agent 最后一条文本（用于解析 DONE:/FAILED:）
    cost_usd: float | None = None
    num_turns: int | None = None
    error: str = ""
    extra: dict = field(default_factory=dict)


LineHook = Callable[[str, dict | None], None]      # (raw_line, parsed_json_or_None)


def _kill_tree(proc: subprocess.Popen) -> None:
    try:
        if _IS_WIN:
            subprocess.run(["taskkill", "/T", "/F", "/PID", str(proc.pid)],
                           capture_output=True, timeout=30)
        else:
            import signal
            try:
                os.killpg(os.getpgid(proc.pid), signal.SIGTERM)
            except Exception:  # noqa: BLE001
                proc.terminate()
            try:
                proc.wait(timeout=10)
            except Exception:  # noqa: BLE001
                try:
                    os.killpg(os.getpgid(proc.pid), signal.SIGKILL)
                except Exception:  # noqa: BLE001
                    proc.kill()
    except Exception as e:  # noqa: BLE001
        logger.warning("结束 Agent 进程树失败 pid=%s: %s", proc.pid, e)


class BaseRunner:
    key = "base"

    def __init__(self, binary: str | None, model: str | None = None):
        self.binary = binary
        self.model = model

    def build_cmd(self, prompt: str, repo_dir: str) -> list[str]:  # pragma: no cover
        raise NotImplementedError

    def parse_line(self, line: str) -> dict | None:
        line = line.strip()
        if not line.startswith("{"):
            return None
        try:
            return json.loads(line)
        except Exception:  # noqa: BLE001
            return None

    def extract_final(self, events: list[dict], result: RunResult) -> None:  # pragma: no cover
        pass

    def humanize(self, ev: dict) -> str | None:  # pragma: no cover
        """把一条事件转成一行可读日志（None=不记录）；子类按各自事件协议实现。"""
        return None

    def run(self, prompt: str, repo_dir: str, env: dict, log_path: str,
            on_line: LineHook | None = None, should_cancel: Callable[[], bool] | None = None,
            timeout_s: int = 2400, on_pid: Callable[[int], None] | None = None) -> RunResult:
        cmd = self.build_cmd(prompt, repo_dir)
        logger.info("启动 Agent [%s]：%s", self.key, " ".join(cmd[:2]) + " …")
        popen_kw: dict = {}
        if not _IS_WIN:
            popen_kw["start_new_session"] = True     # 独立进程组，便于整组结束
        else:
            popen_kw["creationflags"] = subprocess.CREATE_NEW_PROCESS_GROUP
        events: list[dict] = []
        result = RunResult(returncode=-1)
        stream_path = log_path + ".stream.jsonl"     # 原始事件流（JSON 行）；可读日志写 log_path
        with open(log_path, "a", encoding="utf-8", errors="replace") as log, \
                open(stream_path, "a", encoding="utf-8", errors="replace") as stream:
            log.write(f"$ {os.path.basename(cmd[0])} … (cwd={repo_dir})\n")
            log.flush()
            proc = subprocess.Popen(cmd, cwd=repo_dir, env=env, stdin=subprocess.DEVNULL,
                                    stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                                    text=True, encoding="utf-8", errors="replace", bufsize=1, **popen_kw)
            if on_pid:
                on_pid(proc.pid)
            deadline = time.monotonic() + timeout_s
            done = threading.Event()

            def _reader():
                try:
                    for line in proc.stdout:  # type: ignore[union-attr]
                        parsed = self.parse_line(line)
                        if parsed is not None:
                            stream.write(line)
                            stream.flush()
                            events.append(parsed)
                            if len(events) > 5000:      # 只保留最近事件，防内存膨胀
                                del events[:1000]
                            human = self.humanize(parsed)
                            if human:
                                log.write(human.rstrip("\n") + "\n")
                                log.flush()
                        else:
                            log.write(line)
                            log.flush()
                        if on_line:
                            try:
                                on_line(line, parsed)
                            except Exception as e:  # noqa: BLE001
                                logger.debug("on_line 回调异常：%s", e)
                finally:
                    done.set()

            t = threading.Thread(target=_reader, name="pm-agent-reader", daemon=True)
            t.start()
            while not done.wait(timeout=2.0):
                if should_cancel and should_cancel():
                    result.canceled = True
                    log.write("\n[取消] 收到取消请求，结束 Agent 进程\n")
                    _kill_tree(proc)
                    break
                if time.monotonic() > deadline:
                    result.timed_out = True
                    log.write(f"\n[超时] 超过 {timeout_s}s，结束 Agent 进程\n")
                    _kill_tree(proc)
                    break
            done.wait(timeout=30)
            try:
                result.returncode = proc.wait(timeout=60)
            except Exception:  # noqa: BLE001
                _kill_tree(proc)
                result.returncode = proc.poll() if proc.poll() is not None else -9
        try:
            self.extract_final(events, result)
        except Exception as e:  # noqa: BLE001
            logger.debug("解析 Agent 最终输出失败：%s", e)
        return result


def _short(v, n: int = 220) -> str:
    t = v if isinstance(v, str) else json.dumps(v, ensure_ascii=False)
    t = t.replace("\n", " ⏎ ")
    return t if len(t) <= n else t[:n] + "…"


def _tool_summary(name: str, inp: dict) -> str:
    if name == "Bash":
        return f"$ {_short(inp.get('command') or '', 300)}"
    for k in ("file_path", "path", "notebook_path", "pattern", "url", "query"):
        if inp.get(k):
            return f"{name} {_short(inp[k], 200)}"
    return f"{name} {_short(inp, 200)}"


class ClaudeRunner(BaseRunner):
    key = "claude"

    def humanize(self, ev: dict) -> str | None:
        t = ev.get("type")
        if t == "system":
            if ev.get("subtype") == "init":
                return f"[system] 会话开始 model={ev.get('model')} cwd={ev.get('cwd')}"
            return None                       # thinking_tokens 等噪声不记录
        if t == "assistant":
            out = []
            for blk in (ev.get("message") or {}).get("content") or []:
                if not isinstance(blk, dict):
                    continue
                if blk.get("type") == "text" and blk.get("text"):
                    out.append(f"[assistant] {_short(blk['text'], 600)}")
                elif blk.get("type") == "tool_use":
                    out.append(f"[tool] {_tool_summary(blk.get('name') or '?', blk.get('input') or {})}")
            return "\n".join(out) if out else None
        if t == "user":
            out = []
            for blk in (ev.get("message") or {}).get("content") or []:
                if isinstance(blk, dict) and blk.get("type") == "tool_result":
                    c = blk.get("content")
                    if isinstance(c, list):
                        c = " ".join(str(x.get("text", "")) for x in c if isinstance(x, dict))
                    tag = "[tool-error]" if blk.get("is_error") else "[tool-result]"
                    out.append(f"{tag} {_short(c or '', 300)}")
            return "\n".join(out) if out else None
        if t == "result":
            return (f"[result] {'ERROR ' if ev.get('is_error') else ''}turns={ev.get('num_turns')} "
                    f"cost=${ev.get('total_cost_usd')} duration={ev.get('duration_ms')}ms :: "
                    f"{_short(ev.get('result') or '', 800)}")
        return None

    def build_cmd(self, prompt: str, repo_dir: str) -> list[str]:
        s = get_settings()
        cmd = [self.binary or "claude", "-p", prompt, "--output-format", "stream-json", "--verbose",
               "--dangerously-skip-permissions"]
        model = self.model or s.pptmaster_claude_model
        if model:
            cmd += ["--model", model]
        if s.pptmaster_claude_max_budget_usd and s.pptmaster_claude_max_budget_usd > 0:
            cmd += ["--max-budget-usd", str(s.pptmaster_claude_max_budget_usd)]
        return cmd

    def extract_final(self, events: list[dict], result: RunResult) -> None:
        last_text = ""
        for ev in events:
            if ev.get("type") == "assistant":
                for blk in (ev.get("message") or {}).get("content") or []:
                    if isinstance(blk, dict) and blk.get("type") == "text" and blk.get("text"):
                        last_text = blk["text"]
            elif ev.get("type") == "result":
                result.cost_usd = ev.get("total_cost_usd")
                result.num_turns = ev.get("num_turns")
                if ev.get("result"):
                    last_text = str(ev["result"])
                if ev.get("is_error"):
                    result.error = str(ev.get("result") or ev.get("subtype") or "agent error")
        result.final_text = last_text


class CodexRunner(BaseRunner):
    key = "codex"

    def build_cmd(self, prompt: str, repo_dir: str) -> list[str]:
        s = get_settings()
        cmd = [self.binary or "codex", "exec", "--json", "--dangerously-bypass-approvals-and-sandbox",
               "--skip-git-repo-check", "-C", repo_dir]
        model = self.model or s.pptmaster_codex_model
        if model:
            cmd += ["-m", model]
        cmd.append(prompt)
        return cmd

    def humanize(self, ev: dict) -> str | None:
        t = ev.get("type")
        item = ev.get("item") if isinstance(ev.get("item"), dict) else None
        if t == "thread.started":
            return f"[system] 会话开始 thread={ev.get('thread_id')}"
        if item:
            it = item.get("type")
            if it == "command_execution":
                if t == "item.started":
                    return f"[tool] $ {_short(item.get('command') or '', 300)}"
                return f"[tool-result] exit={item.get('exit_code')} {_short(item.get('aggregated_output') or '', 300)}"
            if it == "agent_message" and t == "item.completed":
                return f"[assistant] {_short(item.get('text') or '', 600)}"
            if it == "file_change" and t == "item.completed":
                return f"[file] {_short(item.get('changes') or item, 300)}"
            if it == "reasoning" and t == "item.completed":
                return f"[reasoning] {_short(item.get('text') or '', 200)}"
            return None
        if t == "turn.completed":
            return f"[result] usage={_short(ev.get('usage') or {}, 300)}"
        if t in ("error", "turn.failed"):
            return f"[error] {_short(ev.get('message') or ev.get('error') or ev, 400)}"
        return None

    def extract_final(self, events: list[dict], result: RunResult) -> None:
        last_text = ""
        for ev in events:
            # codex exec --json 事件：{"type":"item.completed","item":{"type":"agent_message","text":...}}
            item = ev.get("item") if isinstance(ev, dict) else None
            if isinstance(item, dict) and item.get("type") == "agent_message" and item.get("text"):
                last_text = str(item["text"])
            elif ev.get("type") == "turn.completed":
                usage = ev.get("usage") or {}
                result.extra["usage"] = usage
            elif ev.get("type") in ("error", "turn.failed"):
                result.error = str(ev.get("message") or ev.get("error") or "codex error")
        result.final_text = last_text


class MockRunner(BaseRunner):
    """不调用任何外部 Agent：直接在 projects/<biz>/exports 下产出占位 PPTX 与两页 SVG 预览。"""
    key = "mock"

    def __init__(self, binary: str | None = None, model: str | None = None,
                 project_dir: str | None = None, pages: int | None = None, title: str = ""):
        super().__init__(binary, model)
        self.project_dir = project_dir
        self.pages = pages or 5
        self.title = title or "ppt-master Mock 演示"

    def build_cmd(self, prompt: str, repo_dir: str) -> list[str]:  # pragma: no cover
        return [sys.executable, "-c", "print('mock')"]

    def run(self, prompt: str, repo_dir: str, env: dict, log_path: str,
            on_line: LineHook | None = None, should_cancel: Callable[[], bool] | None = None,
            timeout_s: int = 2400, on_pid: Callable[[int], None] | None = None) -> RunResult:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        proj = Path(self.project_dir or os.path.join(repo_dir, "projects", "mock"))
        (proj / "exports").mkdir(parents=True, exist_ok=True)
        (proj / "svg_output").mkdir(parents=True, exist_ok=True)
        with open(log_path, "a", encoding="utf-8") as log:
            log.write("[mock] 开始生成占位 PPTX（未调用任何 Agent）\n")
            n = max(1, min(int(self.pages), 60))
            for i in range(1, n + 1):
                if should_cancel and should_cancel():
                    log.write("[mock] 收到取消\n")
                    return RunResult(returncode=1, canceled=True)
                svg = (f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720" width="1280" height="720">'
                       f'<rect width="1280" height="720" fill="#F5F7FA"/><rect x="0" y="0" width="1280" height="88" fill="#1B3A6B"/>'
                       f'<text x="64" y="58" font-size="34" fill="#FFFFFF" font-family="Noto Sans CJK SC, Microsoft YaHei">'
                       f'{self.title} · 第 {i}/{n} 页</text>'
                       f'<text x="64" y="200" font-size="26" fill="#1F2937">Mock 运行器占位页（用于验证提交→轮询→上传→预览链路）</text>'
                       f'<text x="64" y="260" font-size="20" fill="#64748B">提示词长度 {len(prompt)} 字符</text></svg>')
                (proj / "svg_output" / f"P{i:02d}.svg").write_text(svg, encoding="utf-8")
                line = json.dumps({"type": "mock_progress", "page": i, "total": n}, ensure_ascii=False)
                log.write(line + "\n")
                if on_line:
                    on_line(line, {"type": "mock_progress", "page": i, "total": n})
                time.sleep(0.2)
            prs = Presentation()
            prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
            for i in range(1, n + 1):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.9), Inches(1.0))
                tb.text_frame.text = f"{self.title} · 第 {i}/{n} 页"
                tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
                body = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.9), Inches(3.0))
                body.text_frame.text = "Mock 运行器占位内容：真实运行时由 ppt-master 工作流产出原生可编辑 PPTX。"
                body.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
            out = proj / "exports" / f"{proj.name}_mock.pptx"
            prs.save(str(out))
            rel = os.path.relpath(out, repo_dir).replace("\\", "/")
            log.write(f"DONE: {rel}\n")
        return RunResult(returncode=0, final_text=f"DONE: {rel}", num_turns=1, cost_usd=0.0)


def make_runner(agent_key: str, binary: str | None, model: str | None, **mock_kw) -> BaseRunner:
    if agent_key == "claude":
        return ClaudeRunner(binary, model)
    if agent_key == "codex":
        return CodexRunner(binary, model)
    if agent_key == "mock":
        return MockRunner(binary, model, **mock_kw)
    raise ValueError(f"未知 Agent：{agent_key}")


def agent_env(repo_dir: str) -> dict:
    """子进程环境：把后端解释器所在目录前置到 PATH（Agent 调 python3/python 时命中同一 venv），UTF-8 输出。"""
    s = get_settings()
    env = dict(os.environ)
    py = s.pptmaster_python_bin or sys.executable
    py_dir = os.path.dirname(py)
    env["PATH"] = py_dir + os.pathsep + env.get("PATH", "")
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"
    env["PPT_MASTER_HEADLESS"] = "1"
    anthropic_base = env.get("ANTHROPIC_BASE_URL", "").strip()
    if anthropic_base and not anthropic_base.lower().startswith(("http://", "https://")):
        env["ANTHROPIC_BASE_URL"] = "https://" + anthropic_base.lstrip("/")
    # Windows 下 Agent 习惯调用 python3：venv 内没有 python3.exe 时补一个副本（幂等）
    if _IS_WIN and py.lower().endswith("python.exe"):
        py3 = os.path.join(py_dir, "python3.exe")
        if not os.path.exists(py3):
            try:
                shutil.copyfile(py, py3)
            except Exception as e:  # noqa: BLE001
                logger.debug("补 python3.exe 失败（不影响）：%s", e)
    return env
