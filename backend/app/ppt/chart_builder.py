"""图表构建：一律生成 PowerPoint 原生图表（用户可继续编辑数据），失败由上层降级为表格。

数据边界处理（对应文档 §37）：null 剔除、非数值剔除、空数据拒绝、标签去重。
禁止把缺失数据自动当 0。
"""
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Emu, Pt

from app.core.logging import get_logger

logger = get_logger(__name__)

_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


def sanitize_chart_data(data: list) -> list[dict]:
    """清洗图表数据：剔除 null/非数值项，标签去重。清洗后不足 1 个点视为不可用。"""
    cleaned, seen = [], set()
    for item in data or []:
        if not isinstance(item, dict):
            continue
        label, value = str(item.get("label", "")).strip(), item.get("value")
        if not label or label in seen or value is None:
            continue
        try:
            value = float(value)
        except (TypeError, ValueError):
            continue  # 禁止把无法解析的值当 0
        seen.add(label)
        cleaned.append({"label": label, "value": value})
    return cleaned


def add_chart(slide, element: dict, x: Emu, y: Emu, w: Emu, h: Emu, theme,
              font_color: str | None = None, series_color: str | None = None) -> None:
    """失败抛 ValueError，由调用方走降级链路（图表→表格→要点）。
    font_color: 深色背景页传入浅色，保证坐标轴/标签可读。"""
    data = sanitize_chart_data(element.get("data", []))
    if len(data) < 1:
        raise ValueError("图表数据清洗后为空")
    chart_type = _TYPE_MAP.get(element.get("chart_type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = [d["label"] for d in data]
    series_name = element.get("title") or element.get("unit") or "数值"
    chart_data.add_series(series_name, [d["value"] for d in data])

    frame = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data)
    chart = frame.chart
    # 统一样式：小字号、主题色、饼图带图例
    try:
        chart.font.size = Pt(12)
        if font_color:
            chart.font.color.rgb = RGBColor.from_string(font_color)
        if chart_type == XL_CHART_TYPE.PIE:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False
            plot = chart.plots[0]
            series = plot.series[0]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(series_color or theme.primary)
            plot.has_data_labels = True
            plot.data_labels.font.size = Pt(10)
            if font_color:
                plot.data_labels.font.color.rgb = RGBColor.from_string(font_color)
    except Exception as e:  # 样式失败不影响图表主体
        logger.warning("图表样式设置失败(忽略): %s", e)
    logger.info("原生图表构建成功 type=%s 数据点=%d", element.get("chart_type"), len(data))
