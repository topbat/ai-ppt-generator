"""表格构建：原生 PPT 表格，超出容器行数自动截断并标注。"""
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from app.core.logging import get_logger
from app.ppt.fonts import set_run_font

logger = get_logger(__name__)

_ROW_H_EMU = 342900  # 0.375 英寸/行


def add_table(slide, headers: list[str], rows: list[list], x: Emu, y: Emu, w: Emu, h: Emu, theme) -> None:
    if not headers or not isinstance(headers, list):
        raise ValueError("表格缺少表头")
    # 按容器高度限制行数，放不下则截断（QA 报告会记录 Warning）
    max_rows = max(1, int(h / _ROW_H_EMU) - 1)
    body_rows = [r for r in (rows or []) if isinstance(r, list)][:max_rows]
    if len(rows or []) > max_rows:
        logger.warning("表格行数 %d 超出容器，截断为 %d 行", len(rows), max_rows)

    shape = slide.shapes.add_table(len(body_rows) + 1, len(headers), x, y, w,
                                   Emu(_ROW_H_EMU * (len(body_rows) + 1)))
    table = shape.table
    # 表头：主题色底 + 白字加粗
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(theme.primary)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(head)
        set_run_font(run, theme.font_body, 12, "FFFFFF", bold=True)
    # 数据行：隔行浅色
    for r, row in enumerate(body_rows):
        for c in range(len(headers)):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("F5F7FA" if r % 2 else "FFFFFF")
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(row[c]) if c < len(row) else ""
            set_run_font(run, theme.font_body, 11, theme.text)
