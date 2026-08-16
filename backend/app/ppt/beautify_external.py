"""上传 PPT 专业级视觉美化（独立能力，复用 docs/04 视觉优化引擎）。

与生成流水线的 BEAUTIFY 阶段不同：上传的 PPT 没有 Presentation JSON 中间层，
因此这里先用模板解析器的分类器为每页构建"伪 Spec"（版式定型 + 从几何提取信息点），
再套用同一套九维评分 + 确定性美化：

  上传 PPTX → 伪 Spec 构建 → 九维评分（before）
      → 确定性美化：① 纯黑正文 → 深灰 1F2937  ② 左缘锚线聚类吸附  ③ 8pt 网格吸附
      → 复评（after，分数只升不降，否则回退原文件）

铁律：只做吸附级微调与安全的颜色替换，绝不改动内容与版式结构——上传的是
用户资产，"不变差"优先于"变好看"。
"""
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu

from app.core.logging import get_logger
from app.parser.template_parser import _classify_slide
from app.ppt.visual_ops import polish_pptx
from app.ppt.visual_score import score_presentation

logger = get_logger(__name__)

# 分类器产出 → 受控版式枚举的映射（评分器据此套用密度豁免等规则）
_TYPE_MAP = {"content_frame": "title_content", "unknown": "title_content"}


def build_pseudo_spec(pptx_path: str) -> dict:
    """从任意 PPTX 构建评分用伪 Spec：逐页定型 + 正文区文本提取为信息点。"""
    prs = Presentation(pptx_path)
    total = len(prs.slides._sldIdLst)
    slides = []
    for idx, slide in enumerate(prs.slides):
        try:
            stype, _conf, _meta = _classify_slide(slide, idx, total, prs)
        except Exception:
            stype = "title_content"
        stype = _TYPE_MAP.get(stype, stype)
        # 正文区文本（标题带以下）作为信息点，供密度维度评估
        items = []
        for sh in slide.shapes:
            try:
                if (getattr(sh, "has_text_frame", False) and sh.top is not None
                        and Emu(sh.top).inches > 1.3):
                    t = sh.text_frame.text.strip()
                    if t and not t.replace("/", "").replace(" ", "").isdigit():
                        items.append(t[:60])
            except Exception:
                continue
        elements = [{"type": "bullet_group", "items": items[:12]}] if items else []
        slides.append({"page": idx + 1, "type": stype, "title": "", "elements": elements})
    logger.info("伪 Spec 构建完成：%d 页（版式分布 %s）", len(slides),
                {t: [s["type"] for s in slides].count(t) for t in {s["type"] for s in slides}})
    return {"slides": slides}


def _fix_pure_black_text(pptx_path: str) -> int:
    """正文纯黑 (000000) → 深灰 1F2937（视觉规范：正文不用纯黑）。返回修复处数。"""
    prs = Presentation(pptx_path)
    fixed = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if (run.font.color and run.font.color.type is not None
                                and str(run.font.color.rgb) == "000000"):
                            run.font.color.rgb = RGBColor.from_string("1F2937")
                            fixed += 1
                    except Exception:
                        continue
    if fixed:
        prs.save(pptx_path)
    return fixed


def beautify_pptx(in_path: str, out_path: str) -> dict:
    """对上传 PPT 执行专业级视觉美化。返回美化报告（before/after 九维 + 修复明细）。"""
    spec = build_pseudo_spec(in_path)
    before = score_presentation(spec, in_path)
    logger.info("美化前视觉评分：%d 分（%d 页）", before["score"], len(before["pages"]))

    shutil.copyfile(in_path, out_path)
    black_fixed = _fix_pure_black_text(out_path)
    polish = polish_pptx(out_path)
    if black_fixed:
        logger.info("纯黑正文修复：%d 处（000000 → 1F2937）", black_fixed)

    after = score_presentation(spec, out_path)
    if after["score"] < before["score"]:
        # 铁律：复评分数不升反降 → 回退原文件（理论上吸附级微调不会触发，双保险）
        logger.warning("复评 %d < 美化前 %d，回退原文件（不变差原则）",
                       after["score"], before["score"])
        shutil.copyfile(in_path, out_path)
        after = before
        black_fixed = 0
        polish = {"aligned": 0, "spacing_snapped": 0}

    fixes = {"black_text_fixed": black_fixed, "aligned": polish.get("aligned", 0),
             "spacing_snapped": polish.get("spacing_snapped", 0)}
    logger.info("✨ 上传 PPT 美化完成：视觉分 %d → %d｜纯黑修复 %d 处｜对齐吸附 %d 处｜8pt 吸附 %d 处",
                before["score"], after["score"], fixes["black_text_fixed"],
                fixes["aligned"], fixes["spacing_snapped"])
    return {
        "score_before": before["score"],
        "score_after": after["score"],
        "dimensions_before": before["dimensions"],
        "dimensions": after["dimensions"],
        "fixes": fixes,
        "pages": [{"page": p["page"], "score": p["score"], "deductions": p["deductions"]}
                  for p in after["pages"]],
        "total_pages": len(after["pages"]),
    }
