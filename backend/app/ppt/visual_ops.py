"""Fix Ops —— 视觉优化受控 DSL 与确定性调整器（docs/04-VISUAL-OPTIMIZATION.md §5）。

铁律：
- LLM/Vision 只能"建议"枚举内的 op，执行永远是本模块的确定性代码，不会越修越乱；
- spec 级 op 修改 Presentation JSON 后重渲染；pptx 级 op 直接微调渲染产物几何；
- 模板克隆的装饰层只读：pptx 级微调只做 ≤3pt/≤8pt 的吸附级位移，视觉无感但指标归位。
"""
import copy
import re

from pptx import Presentation
from pptx.util import Emu

from app.core.logging import get_logger
from app.schemas.presentation import DENSITY_LIMITS

logger = get_logger(__name__)

# 受控 op 枚举（DSL 全集，枚举外一律拒绝）
FIX_OPS = ("snap_align", "apply_spacing", "font_token", "reduce_palette",
           "upgrade_layout", "rebalance")
# spec 级（改 Presentation JSON 后重渲染） / pptx 级（渲染产物几何微调）
SPEC_OPS = {"upgrade_layout", "rebalance", "font_token"}
PPTX_OPS = {"snap_align", "apply_spacing"}
# reduce_palette：色彩来自模板主题与 Token，spec 无可改字段 → 仅登记留痕（人工/换模板处理）

_NUM_RE = re.compile(r"(\d[\d.,]*\s*(?:%|亿元?|万元?|个|年|家|项|倍|元|人|天|次|吨|km|GW|MW)?)")


# ================= Ops 推导（评分扣分 → 修复指令，确定性映射） =================

def derive_ops(score_result: dict, spec_dict: dict) -> list[dict]:
    """从九维评分结果推导修复指令清单。只产出枚举内 op。"""
    ops: list[dict] = []
    slides = {s["page"]: s for s in spec_dict["slides"]}
    for p in score_result.get("pages", []):
        page, sd = p["page"], slides.get(p["page"])
        if sd is None:
            continue
        by_dim: dict[str, float] = {}
        details: dict[str, list[str]] = {}
        for d in p.get("deductions", []):
            by_dim[d["dim"]] = by_dim.get(d["dim"], 0) + d["points"]
            details.setdefault(d["dim"], []).append(d["detail"])

        if by_dim.get("alignment", 0) >= 3:
            ops.append({"op": "snap_align", "page": page,
                        "detail": "左缘吸附对齐锚线：" + "；".join(details["alignment"])})
        if by_dim.get("spacing", 0) >= 1.5:
            ops.append({"op": "apply_spacing", "page": page,
                        "detail": "间距吸附 8pt 网格：" + "；".join(details["spacing"])})
        if by_dim.get("color", 0) >= 2:
            ops.append({"op": "reduce_palette", "page": page,
                        "detail": "收敛辅色：" + "；".join(details["color"])})
        # 低密度升格：仅普通文字版式（大字报表达本身豁免）
        sparse = any("低密度" in d or "升格" in d for d in details.get("density", []))
        if sparse and sd["type"] in ("title_content", "two_column", "image_text"):
            ops.append({"op": "upgrade_layout", "page": page,
                        "detail": "低密度页升格为大标题/大数字页"})
        # 超预算拥挤 / 截断：精简至预算 → 恢复 Token 字号
        crowded = any("超出版式预算" in d for d in details.get("density", []))
        truncated = any("截断" in d for d in details.get("layout", []))
        if crowded or truncated:
            ops.append({"op": "font_token", "page": page,
                        "detail": "内容精简至预算，恢复 Token 正文字号（拒绝缩字）"})
        # 两栏失衡
        if sd["type"] == "two_column" and _two_column_imbalance(sd) > 2.5:
            ops.append({"op": "rebalance", "page": page,
                        "detail": f"两栏内容失衡（{_two_column_imbalance(sd):.1f} 倍），重新配平"})
    if ops:
        logger.info("Fix Ops 推导完成：%d 条指令 → %s", len(ops),
                    Counter_str([o["op"] for o in ops]))
    else:
        logger.info("Fix Ops 推导完成：无需修复（各维度均在阈值内）")
    return ops


def Counter_str(items: list[str]) -> str:
    from collections import Counter
    return " ".join(f"{k}×{v}" for k, v in Counter(items).items())


def validate_ops(raw_ops: list, total_pages: int) -> list[dict]:
    """校验外部（Vision/LLM）建议的 ops：枚举外/页码非法一律丢弃。"""
    valid = []
    for o in raw_ops or []:
        if not isinstance(o, dict):
            continue
        op, page = o.get("op"), o.get("page")
        if op not in FIX_OPS:
            logger.warning("拒绝枚举外 op：%s（受控 DSL 保护）", op)
            continue
        if not isinstance(page, int) or not (1 <= page <= total_pages):
            logger.warning("拒绝非法页码 op：%s page=%s", op, page)
            continue
        valid.append({"op": op, "page": page, "detail": str(o.get("detail") or "")[:120]})
    return valid


# ================= spec 级调整器（改 JSON → 重渲染） =================

def _two_column_imbalance(sd: dict) -> float:
    """两栏视觉重量失衡倍数（以字符量近似）。"""
    els = [el for el in sd.get("elements", []) if el.get("type") != "image"]
    if len(els) == 2:
        w = [max(1, len(str(el))) for el in els]
        return max(w) / min(w)
    if len(els) == 1 and els[0].get("type") == "bullet_group":
        items = [str(i) for i in els[0].get("items") or []]
        if len(items) >= 2:
            mid = (len(items) + 1) // 2
            a, b = sum(map(len, items[:mid])), sum(map(len, items[mid:]))
            return max(a, b) / max(1, min(a, b))
    return 1.0


def _apply_upgrade_layout(sd: dict) -> bool:
    """低密度页升格：含数字 → 大数字页；否则 → 大标题页（quote 版式 40pt 居中）。"""
    texts: list[str] = []
    for el in sd.get("elements", []):
        if el.get("type") == "bullet_group":
            texts += [str(i) for i in el.get("items") or []]
        elif el.get("type") == "paragraph" and el.get("text"):
            texts.append(str(el["text"]))
    numbers = []
    for t in texts:
        m = _NUM_RE.search(t)
        if m and any(c.isdigit() for c in m.group(1)):
            label = _NUM_RE.sub("", t).strip("，。、：:；; ")[:12] or sd.get("title", "")[:12]
            numbers.append({"value": m.group(1).strip(), "label": label})
    if numbers:
        sd["type"] = "key_number"
        sd["elements"] = [{"type": "key_number", "items": numbers[:3]}]
        logger.info("第 %s 页升格为大数字页（%d 个关键数字成为视觉主体）",
                    sd["page"], len(numbers[:3]))
        return True
    big_text = sd.get("key_message") or (texts[0] if texts else "") or sd.get("title", "")
    if not big_text:
        return False
    sd["type"] = "quote"
    sd["elements"] = [{"type": "quote", "text": str(big_text)[:60], "author": None}]
    logger.info("第 %s 页升格为大标题页（低密度内容改为居中大字表达）", sd["page"])
    return True


def _apply_rebalance(sd: dict) -> bool:
    """两栏配平：单要点组按字符量贪心均分为两组。"""
    els = sd.get("elements", [])
    bg = next((el for el in els if el.get("type") == "bullet_group"), None)
    if bg is None or len(bg.get("items") or []) < 3:
        return False
    items = [str(i) for i in bg["items"]]
    left: list[str] = []
    right: list[str] = []
    lw = rw = 0
    for it in sorted(items, key=len, reverse=True):
        if lw <= rw:
            left.append(it)
            lw += len(it)
        else:
            right.append(it)
            rw += len(it)
    # 保持原有阅读顺序（组内按原顺序排列）
    left.sort(key=items.index)
    right.sort(key=items.index)
    sd["elements"] = ([{"type": "bullet_group", "items": left},
                       {"type": "bullet_group", "items": right}]
                      + [el for el in els if el is not bg and el.get("type") != "bullet_group"])
    logger.info("第 %s 页两栏重新配平：%d+%d 条（字符量 %d/%d）",
                sd["page"], len(left), len(right), lw, rw)
    return True


def _apply_font_token(sd: dict, density: str) -> bool:
    """内容精简至密度预算：减条目/截长句 → 渲染时自动恢复 Token 字号。"""
    limits = DENSITY_LIMITS.get(density, DENSITY_LIMITS["medium"])
    changed = False
    for el in sd.get("elements", []):
        if el.get("type") == "bullet_group":
            items = [str(i) for i in el.get("items") or []]
            trimmed = [i if len(i) <= limits["bullet"] else i[:limits["bullet"] - 1] + "…"
                       for i in items[:limits["max_items"]]]
            if trimmed != items:
                el["items"] = trimmed
                changed = True
        elif el.get("type") == "cards":
            for it in el.get("items") or []:
                desc = str(it.get("desc") or "")
                if len(desc) > limits["card_desc"]:
                    it["desc"] = desc[:limits["card_desc"] - 1] + "…"
                    changed = True
    if changed:
        logger.info("第 %s 页内容精简至密度预算（字号回归 Token，拒绝缩字适配）", sd["page"])
    return changed


def apply_spec_ops(spec_dict: dict, ops: list[dict], density: str = "medium"
                   ) -> tuple[dict, list[dict], list[dict]]:
    """应用 spec 级 ops。返回 (新 spec dict, 已应用, 跳过)。输入不被修改。"""
    new_spec = copy.deepcopy(spec_dict)
    slides = {s["page"]: s for s in new_spec["slides"]}
    applied, skipped = [], []
    for op in ops:
        sd = slides.get(op["page"])
        if sd is None or op["op"] not in SPEC_OPS:
            if op["op"] not in PPTX_OPS:  # pptx 级另行处理，不算跳过
                skipped.append(op)
            continue
        try:
            ok = {"upgrade_layout": _apply_upgrade_layout,
                  "rebalance": _apply_rebalance,
                  }.get(op["op"], lambda s: _apply_font_token(s, density))(sd)
            (applied if ok else skipped).append(op)
        except Exception as e:
            logger.warning("op %s 应用失败（保留原页）page=%s: %s", op["op"], op["page"], e)
            skipped.append(op)
    for op in ops:
        if op["op"] == "reduce_palette":
            logger.info("第 %s 页 reduce_palette 登记留痕（色彩源于模板主题，规则层不改色）",
                        op["page"])
    return new_spec, applied, skipped


# ================= pptx 级微调（渲染产物几何吸附，全模式渲染期执行） =================

def polish_pptx(pptx_path: str) -> dict:
    """渲染期一次性正确：对齐锚线聚类吸附（≤8pt）+ 8pt 网格吸附（≤3pt）。

    只微调带文字的形状、只做吸附级小位移——模板装饰层视觉无感，
    但 Alignment / Spacing 维度指标确定性归位。"""
    prs = Presentation(pptx_path)
    aligned = snapped = 0
    for slide in prs.slides:
        text_shapes = [sh for sh in slide.shapes
                       if getattr(sh, "has_text_frame", False)
                       and sh.text_frame.text.strip() and sh.left is not None]
        # 1) 左缘聚类吸附：只把"游离"形状（未与他人成簇）拉向最近的主锚线（容差 8pt）；
        #    已成簇（≥2 成员）的形状保持不动，避免相邻锚线互相拉扯
        lefts = [round(Emu(sh.left).inches, 3) for sh in text_shapes]
        clusters = _cluster_anchors(lefts)
        majors = [a for a, cnt in clusters if cnt >= 2]
        strays = {a for a, cnt in clusters if cnt == 1}
        for sh in text_shapes:
            cur = Emu(sh.left).inches
            if not any(abs(cur - s) <= 0.03 for s in strays):
                continue
            near = [a for a in majors if abs(cur - a) <= 8 / 72]
            if near:
                target = min(near, key=lambda a: abs(cur - a))
                sh.left = Emu(int(target * 914400))
                aligned += 1
        # 2) 垂直 8pt 网格吸附（容差 3pt，不产生可感知位移）
        for sh in text_shapes:
            top_pt = Emu(sh.top).inches * 72
            target = round(top_pt / 8) * 8
            if 0.1 < abs(top_pt - target) <= 3.0:
                sh.top = Emu(int(target / 72 * 914400))
                snapped += 1
    if aligned or snapped:
        prs.save(pptx_path)
        logger.info("渲染微调完成：左缘吸附 %d 处、8pt 网格吸附 %d 处（一次性正确）",
                    aligned, snapped)
    return {"aligned": aligned, "spacing_snapped": snapped}


def _cluster_anchors(lefts: list[float]) -> list[tuple[float, int]]:
    """左缘值聚类 → [(锚线, 成员数)]，容差 0.03 英寸。"""
    clusters: list[list[float]] = []
    for x in sorted(lefts):
        if clusters and x - clusters[-1][0] <= 0.03:
            clusters[-1].append(x)
        else:
            clusters.append([x])
    return [(min(c), len(c)) for c in sorted(clusters, key=len, reverse=True)]
