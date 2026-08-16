"""PPT 视觉评分模型（docs/04-VISUAL-OPTIMIZATION.md §4）。

九维 100 分制规则引擎：逐页评分 + 整册汇总，全部确定性计算（毫秒级、零 LLM 成本）。
输入 = Presentation JSON（内容结构）+ 渲染产物 PPTX（真实几何/字号/颜色），
输出 = {score, dimensions, pages:[{page, score, dimensions, deductions}]}。

与 quality_score（生成质量分：对不对）并列，visual_score 只回答"好不好看"。
"""
from collections import Counter

from pptx import Presentation
from pptx.util import Emu

from app.core.logging import get_logger
from app.ppt import design_tokens as dt

logger = get_logger(__name__)

# 九维满分（合计 100）
DIM_MAX = {
    "layout": 20, "alignment": 15, "typography": 15, "spacing": 10, "color": 10,
    "hierarchy": 10, "density": 10, "image": 5, "consistency": 5,
}
DIM_NAMES = {
    "layout": "布局", "alignment": "对齐", "typography": "字体", "spacing": "间距",
    "color": "色彩", "hierarchy": "层级", "density": "密度", "image": "图片",
    "consistency": "一致性",
}

# 全幅装饰页（封面/章节/尾页）：几何维度豁免（构图自由，不按正文网格评判）
_FULL_BLEED_TYPES = {"cover", "section", "ending"}
# 低密度豁免版式（本身就是"大字报"表达，稀疏是特性不是缺陷）
_SPARSE_EXEMPT_TYPES = _FULL_BLEED_TYPES | {"quote", "key_number", "toc", "summary"}

# 密度预算（信息点数区间，思想文档：普通 3~5 / 卡片 3~4 / 架构 5~9）
_DENSITY_BUDGET = {
    "title_content": (3, 6), "two_column": (3, 8), "image_text": (3, 6), "summary": (2, 6),
    "three_column": (3, 4), "three_cards": (3, 4), "four_cards": (3, 5),
    "architecture": (4, 16), "timeline": (3, 6), "process": (3, 6),
    "comparison": (4, 10), "table": (2, 10),
    "bar_chart": (2, 10), "line_chart": (2, 12), "pie_chart": (2, 8),
}

_EMU_IN = 914400.0
_SNAP_TOL_PT = 3.0        # 8pt 网格命中容差
_ALIGN_TOL_IN = 0.08      # 对齐锚线聚类容差（英寸）


# ================= 几何采集 =================

def _shape_info(shape) -> dict | None:
    """采集单个形状的几何与文本属性；无位置信息的形状跳过。"""
    try:
        if shape.left is None or shape.top is None:
            return None
        info = {
            "left": Emu(shape.left).inches, "top": Emu(shape.top).inches,
            "width": Emu(shape.width or 0).inches, "height": Emu(shape.height or 0).inches,
            "is_picture": shape.shape_type == 13,  # MSO_SHAPE_TYPE.PICTURE
            "runs": [], "text": "",
        }
    except Exception:
        return None
    if getattr(shape, "has_text_frame", False):
        texts = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                t = run.text or ""
                if not t.strip():
                    continue
                texts.append(t)
                color = None
                try:
                    if run.font.color and run.font.color.type is not None:
                        color = str(run.font.color.rgb)
                except Exception:
                    color = None
                info["runs"].append({
                    "size": run.font.size.pt if run.font.size else None,
                    "bold": bool(run.font.bold), "color": color,
                    "font": run.font.name, "chars": len(t),
                })
        info["text"] = "".join(texts)
    return info


def collect_slide_geometry(slide) -> list[dict]:
    """采集一页内全部形状信息（组合形状降级为整体外框）。"""
    infos = []
    for shape in slide.shapes:
        info = _shape_info(shape)
        if info is not None:
            infos.append(info)
    return infos


# ================= 颜色工具 =================

def _hex_rgb(hex_color: str) -> tuple[int, int, int]:
    return int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)


def _is_neutral(hex_color: str) -> bool:
    """中性色判定：灰阶（饱和度低）或接近纯白/纯黑。"""
    r, g, b = _hex_rgb(hex_color)
    mx, mn = max(r, g, b), min(r, g, b)
    lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return (mx - mn) < 30 or lum > 0.93 or lum < 0.10


def _hue_bucket(hex_color: str) -> int:
    """色相量化到 30° 桶，同桶视作同一色系。"""
    r, g, b = _hex_rgb(hex_color)
    mx, mn = max(r, g, b), min(r, g, b)
    if mx == mn:
        return -1
    d = mx - mn
    if mx == r:
        h = (60 * ((g - b) / d) + 360) % 360
    elif mx == g:
        h = 60 * ((b - r) / d) + 120
    else:
        h = 60 * ((r - g) / d) + 240
    return int(h // 30)


# ================= 信息点计数（来自 Presentation JSON） =================

def info_points(sd: dict) -> int:
    """一页的信息点数（密度维度的输入）。"""
    n = 0
    for el in sd.get("elements", []):
        t = el.get("type")
        if t in ("bullet_group", "cards", "key_number", "timeline"):
            n += len(el.get("items") or [])
        elif t == "process":
            n += len(el.get("steps") or [])
        elif t == "architecture":
            n += len(el.get("layers") or [])
            n += sum(len(la.get("items") or []) for la in (el.get("layers") or []))
        elif t == "comparison":
            n += len((el.get("left") or {}).get("items") or [])
            n += len((el.get("right") or {}).get("items") or [])
        elif t == "table":
            n += len(el.get("rows") or [])
        elif t == "chart":
            n += len(el.get("data") or [])
        elif t in ("quote", "paragraph"):
            n += 1
        elif t == "image":
            n += 1
    return n


# ================= 单页九维评分 =================

def _ded(deductions: list, dim: str, points: float, detail: str) -> float:
    deductions.append({"dim": dim, "points": round(points, 1), "detail": detail})
    return points


def _score_layout(sd, shapes, page_w, page_h, notes, deductions) -> float:
    score = DIM_MAX["layout"]
    # 1) 元素外溢画布（硬伤）
    overflow = [s for s in shapes
                if s["left"] < -0.05 or s["top"] < -0.05
                or s["left"] + s["width"] > page_w + 0.06
                or s["top"] + s["height"] > page_h + 0.06]
    if overflow:
        score -= _ded(deductions, "layout", min(10, 5 * len(overflow)),
                      f"{len(overflow)} 个元素超出画布边界")
    # 2) 渲染期问题（截断/降级）
    truncated = [n for n in notes if n.get("issue") == "text_truncated"]
    if truncated:
        score -= _ded(deductions, "layout", min(6, 3 * len(truncated)),
                      f"{len(truncated)} 处文本被截断（容器不足）")
    if any(n.get("issue") == "page_degraded" for n in notes):
        score -= _ded(deductions, "layout", 8, "页面曾降级绘制")
    # 3) 正文页空洞：正文区（标题带以下）无任何文本
    if sd["type"] not in _FULL_BLEED_TYPES:
        body_texts = [s for s in shapes if s["runs"] and s["top"] > 1.2]
        if not body_texts and not any(s["is_picture"] for s in shapes):
            score -= _ded(deductions, "layout", 6, "正文区无内容（页面空洞）")
    return max(0, score)


def _score_alignment(sd, shapes, page_w, deductions) -> float:
    score = DIM_MAX["alignment"]
    if sd["type"] in _FULL_BLEED_TYPES:
        return score
    texted = [s for s in shapes if s["runs"]]
    if len(texted) <= 1:
        return score
    lefts, centered = [], 0
    for s in texted:
        cx = s["left"] + s["width"] / 2
        if abs(cx - page_w / 2) < 0.2:   # 居中构图的形状不参与左缘锚线统计
            centered += 1
        else:
            lefts.append(round(s["left"], 3))
    if len(lefts) <= 1:
        return score
    # 左缘聚类：取前 2 个主锚线，容差外的即"游离元素"
    clusters: list[list[float]] = []
    for x in sorted(lefts):
        if clusters and x - clusters[-1][0] <= _ALIGN_TOL_IN:
            clusters[-1].append(x)
        else:
            clusters.append([x])
    clusters.sort(key=len, reverse=True)
    anchors = [sum(c) / len(c) for c in clusters[:2]]
    stray = sum(1 for x in lefts if all(abs(x - a) > _ALIGN_TOL_IN for a in anchors))
    if stray:
        score -= _ded(deductions, "alignment", min(9, 3 * stray),
                      f"{stray} 个元素左缘偏离对齐锚线")
    return max(0, score)


def _score_typography(sd, shapes, deductions) -> float:
    score = DIM_MAX["typography"]
    sizes = [r["size"] for s in shapes for r in s["runs"] if r["size"]]
    if not sizes:
        return score
    tiny = sorted({s for s in sizes if s < 10})
    if tiny:
        score -= _ded(deductions, "typography", min(4.5, 1.5 * len(tiny)),
                      f"存在过小字号 {tiny}pt（可读性风险）")
    distinct = len({round(s) for s in sizes})
    if distinct > 6:
        score -= _ded(deductions, "typography", 5 if distinct > 8 else 3,
                      f"单页字号种类过多（{distinct} 种）")
    # 标题主导性：最大字号应出现在标题带（全幅页豁免）
    if sd["type"] not in _FULL_BLEED_TYPES:
        mx = max(sizes)
        top_max = max((r["size"] for s in shapes if s["top"] < 1.4
                       for r in s["runs"] if r["size"]), default=0)
        if top_max and mx > top_max + 2:
            score -= _ded(deductions, "typography", 3,
                          f"正文区字号({mx}pt)压过标题({top_max}pt)，层级倒挂")
    return max(0, score)


def _score_spacing(sd, shapes, deductions) -> float:
    score = DIM_MAX["spacing"]
    if sd["type"] in _FULL_BLEED_TYPES:
        return score
    tops = [s["top"] for s in shapes if s["runs"]]
    if len(tops) < 2:
        return score
    miss = 0
    for t in tops:
        pt = t * 72
        if abs(pt - round(pt / 8) * 8) > _SNAP_TOL_PT:
            miss += 1
    rate = 1 - miss / len(tops)
    if rate < 0.999:
        lost = round(DIM_MAX["spacing"] * (1 - rate) * 0.7, 1)
        if lost >= 0.5:
            score -= _ded(deductions, "spacing", lost,
                          f"{miss}/{len(tops)} 个元素未吸附 8pt 网格")
    return max(0, score)


def _score_color(sd, shapes, deductions) -> float:
    score = DIM_MAX["color"]
    colors = [r["color"] for s in shapes for r in s["runs"] if r["color"]]
    if not colors:
        return score
    hues = {_hue_bucket(c) for c in colors if not _is_neutral(c)}
    hues.discard(-1)
    if len(hues) > 3:
        score -= _ded(deductions, "color", min(6, 2 * (len(hues) - 3)),
                      f"色相过多（{len(hues)} 个色系 > 主1+辅2）")
    if sd["type"] not in _FULL_BLEED_TYPES and "000000" in colors:
        score -= _ded(deductions, "color", 2, "正文使用纯黑（应使用深灰 1F2937）")
    return max(0, score)


def _score_hierarchy(sd, shapes, deductions) -> float:
    score = DIM_MAX["hierarchy"]
    if sd["type"] in _FULL_BLEED_TYPES:
        return score
    runs = [r for s in shapes for r in s["runs"] if r["size"]]
    if len(runs) < 3:
        return score
    strengths = [round(r["size"] * (1.25 if r["bold"] else 1.0)) for r in runs]
    levels = sorted(set(strengths), reverse=True)
    if len(levels) == 1 and len(runs) >= 4:
        score -= _ded(deductions, "hierarchy", 4, "全页视觉强度扁平（无层级差异）")
    elif len(levels) > 5:
        score -= _ded(deductions, "hierarchy", 3, f"视觉层级过多（{len(levels)} 级）")
    else:
        top_n = strengths.count(levels[0])
        if top_n >= 3 and top_n / len(strengths) > 0.5:
            score -= _ded(deductions, "hierarchy", 3, "最高强度元素过多（强调过度=没有强调）")
    return max(0, score)


def _score_density(sd, deductions) -> float:
    score = DIM_MAX["density"]
    stype = sd["type"]
    n = info_points(sd)
    if stype in _SPARSE_EXEMPT_TYPES:
        return score
    lo, hi = _DENSITY_BUDGET.get(stype, (3, 6))
    if n == 0:
        score -= _ded(deductions, "density", 6, "正文无信息点")
    elif n <= 2 and stype in ("title_content", "two_column", "image_text"):
        score -= _ded(deductions, "density", 4,
                      f"信息点仅 {n} 个（低密度，建议升格为大标题/大数字页）")
    elif n < lo:
        score -= _ded(deductions, "density", 2, f"信息点 {n} 低于版式预算下限 {lo}")
    elif n > hi:
        score -= _ded(deductions, "density", min(6, 1.5 * (n - hi)),
                      f"信息点 {n} 超出版式预算上限 {hi}（拥挤）")
    return max(0, score)


def _score_image(sd, shapes, page_w, page_h, deductions) -> float:
    score = DIM_MAX["image"]
    pics = [s for s in shapes if s["is_picture"]]
    if not pics or sd["type"] in _FULL_BLEED_TYPES:
        return score
    area = sum(p["width"] * p["height"] for p in pics) / (page_w * page_h)
    if area > 0.65:
        score -= _ded(deductions, "image", 2, f"图片占比 {area:.0%} 过大")
    # 图片与文本叠压检测
    for p in pics:
        for s in shapes:
            if not s["runs"]:
                continue
            ox = max(0.0, min(p["left"] + p["width"], s["left"] + s["width"]) - max(p["left"], s["left"]))
            oy = max(0.0, min(p["top"] + p["height"], s["top"] + s["height"]) - max(p["top"], s["top"]))
            s_area = max(1e-6, s["width"] * s["height"])
            if ox * oy / s_area > 0.3:
                score -= _ded(deductions, "image", 2, "图片与文本叠压")
                return max(0, score)
    return max(0, score)


def _score_consistency(page_infos: list[dict]) -> tuple[float, list[dict]]:
    """整册一致性（跨页维度）：标题锚线、字体统一、页码存在率。"""
    score, deductions = DIM_MAX["consistency"], []
    content = [p for p in page_infos if p["type"] not in _FULL_BLEED_TYPES]
    if len(content) >= 3:
        # 1) 标题左缘一致性
        title_lefts = [p["title_left"] for p in content if p.get("title_left") is not None]
        if len(title_lefts) >= 3:
            spread = max(title_lefts) - min(title_lefts)
            if spread > 0.3:
                score -= _ded(deductions, "consistency", 2,
                              f"各页标题左缘漂移 {spread:.2f}\"（锚线不一致）")
        # 2) 字体统一性
        fonts = Counter(f for p in content for f in p["fonts"] if f)
        if fonts:
            dominant = fonts.most_common(1)[0][1] / sum(fonts.values())
            if dominant < 0.7:
                score -= _ded(deductions, "consistency", 1.5,
                              f"字体不统一（主字体占比 {dominant:.0%}）")
        # 3) 页码存在率
        with_no = sum(1 for p in content if p["has_page_no"])
        if with_no / len(content) < 0.6:
            score -= _ded(deductions, "consistency", 1,
                          f"页码缺失（{with_no}/{len(content)} 页有页码）")
    return max(0, score), deductions


# ================= 主入口 =================

def score_presentation(spec, pptx_path: str, notes: list[dict] | None = None) -> dict:
    """spec: PresentationSpec 或等价 dict；pptx_path: 渲染产物；notes: 渲染备注。"""
    sd_list = spec.model_dump()["slides"] if hasattr(spec, "model_dump") else spec["slides"]
    sd_list = sorted(sd_list, key=lambda s: s["page"])
    notes_by_page: dict[int, list] = {}
    for n in (notes or []):
        notes_by_page.setdefault(n.get("page"), []).append(n)

    prs = Presentation(pptx_path)
    page_w = Emu(prs.slide_width).inches if prs.slide_width else dt.CANVAS_W
    page_h = Emu(prs.slide_height).inches if prs.slide_height else dt.CANVAS_H
    slides = list(prs.slides)

    pages, page_infos = [], []
    for idx, sd in enumerate(sd_list):
        shapes = collect_slide_geometry(slides[idx]) if idx < len(slides) else []
        page_notes = notes_by_page.get(sd["page"], [])
        deductions: list[dict] = []
        dims = {
            "layout": _score_layout(sd, shapes, page_w, page_h, page_notes, deductions),
            "alignment": _score_alignment(sd, shapes, page_w, deductions),
            "typography": _score_typography(sd, shapes, deductions),
            "spacing": _score_spacing(sd, shapes, deductions),
            "color": _score_color(sd, shapes, deductions),
            "hierarchy": _score_hierarchy(sd, shapes, deductions),
            "density": _score_density(sd, deductions),
            "image": _score_image(sd, shapes, page_w, page_h, deductions),
        }
        # 一致性为整册维度，页级先记满分占位，汇总后统一回填
        pages.append({"page": sd["page"], "type": sd["type"], "dims": dims,
                      "deductions": deductions})
        # 跨页一致性所需的页级采样
        title_shape = next((s for s in shapes if s["runs"] and s["top"] < 1.4
                            and any(r["size"] and r["size"] >= 15 for r in s["runs"])), None)
        page_infos.append({
            "type": sd["type"],
            "title_left": round(title_shape["left"], 3) if title_shape else None,
            "fonts": [r["font"] for s in shapes for r in s["runs"]],
            "has_page_no": any(
                s["top"] > page_h - 0.75 and s["left"] > page_w * 0.65
                and s["text"].strip().isdigit() for s in shapes),
        })

    consistency, cons_deductions = _score_consistency(page_infos)

    dim_totals: dict[str, float] = {k: 0.0 for k in DIM_MAX}
    for p in pages:
        p["dims"]["consistency"] = consistency
        p["score"] = round(sum(p["dims"].values()))
        for k, v in p["dims"].items():
            dim_totals[k] += v
    n_pages = max(1, len(pages))
    dimensions = {k: {"score": round(dim_totals[k] / n_pages, 1), "max": DIM_MAX[k],
                      "name": DIM_NAMES[k]} for k in DIM_MAX}
    total = round(sum(d["score"] for d in dimensions.values()))

    result = {
        "score": total,
        "dimensions": dimensions,
        "deck_deductions": cons_deductions,
        "pages": [{"page": p["page"], "type": p["type"], "score": p["score"],
                   "dimensions": {k: round(v, 1) for k, v in p["dims"].items()},
                   "deductions": p["deductions"]} for p in pages],
    }
    worst = min(pages, key=lambda p: p["score"], default=None)
    logger.info("视觉评分完成：整册 %d 分（%s）| 最低页 P%s=%s 分 | 扣分项 %d 处",
                total,
                " ".join(f"{DIM_NAMES[k]}{dimensions[k]['score']}" for k in DIM_MAX),
                worst["page"] if worst else "-", worst["score"] if worst else "-",
                sum(len(p["deductions"]) for p in pages) + len(cons_deductions))
    return result
