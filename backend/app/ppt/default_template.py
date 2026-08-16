"""系统默认模板：无需用户上传即可生成，程序化构建（首次启动时注册入库）。"""
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

_PRIMARY = "1B3A6B"
_ACCENT = "2F80ED"


def build_default_template() -> bytes:
    """构建一份简洁商务风模板（封面/目录/章节/正文/卡片/总结 6 页示意）。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide(bg: str | None = None):
        slide = prs.slides.add_slide(blank)
        if bg:
            rect = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor.from_string(bg)
            rect.line.fill.background()
        return slide

    def add_text(slide, text, x, y, w, h, size, color=_PRIMARY, bold=True):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = RGBColor.from_string(color)
        return tb

    # 1 封面
    s = add_slide(_PRIMARY)
    add_text(s, "演示文稿标题", 1.2, 2.8, 11, 1.2, 40, "FFFFFF")
    add_text(s, "副标题说明", 1.2, 4.1, 11, 0.8, 20, "B9CBE8", bold=False)
    # 2 目录
    s = add_slide()
    add_text(s, "目录", 0.8, 0.5, 4, 0.8, 30)
    for i in range(4):
        add_text(s, f"{i + 1}、章节标题", 1.5, 1.8 + i * 1.1, 8, 0.7, 20, "333333", bold=False)
    # 3 章节页
    s = add_slide(_PRIMARY)
    add_text(s, "01", 1.0, 2.2, 3, 1.5, 60, _ACCENT)
    add_text(s, "章节标题", 1.0, 3.9, 8, 1.0, 32, "FFFFFF")
    # 4 正文页
    s = add_slide()
    add_text(s, "页面标题", 0.8, 0.5, 8, 0.8, 28)
    add_text(s, "正文要点内容示意", 0.8, 1.8, 11.5, 4.5, 16, "333333", bold=False)
    # 5 三卡片页
    s = add_slide()
    add_text(s, "页面标题", 0.8, 0.5, 8, 0.8, 28)
    for i in range(3):
        card = s.shapes.add_shape(1, Inches(0.8 + i * 4.1), Inches(2.0), Inches(3.7), Inches(3.6))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("F5F7FA")
        card.line.color.rgb = RGBColor.from_string("D9E2F1")
        add_text(s, f"卡片{i + 1}", 1.1 + i * 4.1, 2.3, 3.1, 0.6, 18)
    # 6 总结页
    s = add_slide(_PRIMARY)
    add_text(s, "总结", 1.2, 3.0, 10, 1.2, 36, "FFFFFF")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
