"""模板页克隆引擎：在同一 Presentation 内复制模板页并替换文字。

这是"保留模板风格"的核心能力（文档 §32 情况 A 的升级实现）：
- 封面/目录/尾页：整页克隆（背景、装饰、图片、字体全部保留），只替换文字；
- 正文页：克隆"内容框架页"作为底版，标题写入原标题形状，内容绘制在正文区。

实现要点：
- 形状树 deepcopy + 关系(rels)重建 + rId 重映射（图片等引用不丢失）；
- 文字替换保留原 run 的字体样式（改 text 不动 rPr）；
- 输出顺序 = 克隆顺序，最后删除模板原始页。
"""
import copy
import re

from pptx.oxml.ns import qn

from app.core.logging import get_logger

logger = get_logger(__name__)

_RELTYPE_SLIDE_LAYOUT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
_RELTYPE_NOTES = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"


def _copy_shapes_with_rels(new_slide, elements, src_part) -> list:
    """把形状元素深拷贝进 new_slide，并重建来自 src_part 的关系。

    rId 重映射只作用于本批新增节点（两阶段替换避免交叉覆盖），
    支持同一页先后从版式部件与页面部件复制两批形状。"""
    appended = []
    for el in elements:
        c = copy.deepcopy(el)
        new_slide.shapes._spTree.append(c)
        appended.append(c)

    mapping: dict[str, str] = {}
    for rId, rel in src_part.rels.items():
        if rel.reltype in (_RELTYPE_SLIDE_LAYOUT, _RELTYPE_NOTES):
            continue
        try:
            if rel.is_external:
                new_rId = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            mapping[rId] = new_rId
        except Exception as e:
            logger.warning("克隆页关系复制失败(忽略该引用) rId=%s: %s", rId, e)

    if mapping:
        tokens = {old: f"__TMP_{i}__" for i, old in enumerate(mapping)}
        for root in appended:
            for el in root.iter():
                for attr, val in list(el.attrib.items()):
                    if val in tokens and (attr.endswith("}embed") or attr.endswith("}link")
                                          or attr.endswith("}id") or attr.endswith("}pict")):
                        el.set(attr, tokens[val])
        reverse = {tokens[old]: new for old, new in mapping.items()}
        for root in appended:
            for el in root.iter():
                for attr, val in list(el.attrib.items()):
                    if val in reverse:
                        el.set(attr, reverse[val])
    return appended


def _copy_bg(new_slide, source_slide, include_layout: bool = False) -> None:
    """复制背景（页面级 p:cSld/p:bg，可选回退到版式背景）。"""
    src_bg = source_slide._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if src_bg is None and include_layout:
        src_bg = source_slide.slide_layout._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if src_bg is not None:
        new_cSld = new_slide._element.find(qn("p:cSld"))
        new_cSld.insert(0, copy.deepcopy(src_bg))


def clone_slide(prs, source_slide):
    """克隆一页（含背景/装饰/图片），追加到演示文稿末尾，返回新页。"""
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shp in list(new_slide.shapes):  # 清掉 add_slide 从版式带入的占位符
        shp._element.getparent().remove(shp._element)
    _copy_bg(new_slide, source_slide)
    _copy_shapes_with_rels(new_slide, [s._element for s in source_slide.shapes],
                           source_slide.part)
    return new_slide


def _layout_shape_contaminated(shp, page_area: int, page_h_emu: int) -> bool:
    """判断版式形状是否为污染物（示例图/感谢艺术字/演示正文），拍平时丢弃。"""
    try:
        text = " ".join(s.text_frame.text for s in iter_text_shapes([shp])).strip()
        if text and _ENDING_TEXT_RE.search(text):
            return True  # 版式携带"谢谢"类艺术字/文本
        area = (shp.width or 0) * (shp.height or 0)
        frac = area / max(page_area, 1)
        has_pic = shp.shape_type == 13 or (
            shp.shape_type == 6 and any(s.shape_type == 13 for s in shp.shapes))
        if has_pic and 0.06 < frac < 0.9:
            return True  # 版式内中等尺寸示例图（全幅背景图保留）
        if text and len(text.replace(" ", "")) > 2 and shp.top is not None and shp.height:
            top_f = shp.top / max(page_h_emu, 1)
            bot_f = (shp.top + shp.height) / max(page_h_emu, 1)
            if not (bot_f < 0.16 or top_f > 0.84):
                return True  # 非页眉页脚区的版式演示文字
    except Exception:
        return False
    return False


def clone_slide_flat(prs, source_slide, blank_layout):
    """拍平克隆：版式(layout)可见形状垫底拍进页面 + 页面形状置顶，版式改指空白版式。

    解决"版式层携带示例照片/'谢谢'艺术字随克隆污染所有正文页"的问题——
    拍平后这些形状成为页面形状，统一接受后续清理规则（正文区/标题带/尾字清理）。"""
    new_slide = prs.slides.add_slide(blank_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    _copy_bg(new_slide, source_slide, include_layout=True)

    layout = source_slide.slide_layout
    page_area = int(prs.slide_width or 1) * int(prs.slide_height or 1)
    page_h = int(prs.slide_height or 1)
    lay_elems = []
    for shp in layout.shapes:
        try:
            if getattr(shp, "is_placeholder", False):
                continue  # 版式占位符只是样式原型，不参与页面渲染
            if _layout_shape_contaminated(shp, page_area, page_h):
                continue
            lay_elems.append(shp._element)
        except Exception:
            continue
    _copy_shapes_with_rels(new_slide, lay_elems, layout.part)          # 版式装饰垫底
    _copy_shapes_with_rels(new_slide, [s._element for s in source_slide.shapes],
                           source_slide.part)                           # 页面形状在上
    return new_slide


def delete_slide(prs, index: int) -> None:
    """按索引删除页（从 sldIdLst 移除并断开关系）。"""
    sldIdLst = prs.slides._sldIdLst
    slide_id = list(sldIdLst)[index]
    rId = slide_id.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(slide_id)


def delete_leading_slides(prs, count: int) -> None:
    """删除最前面的 count 页（即模板原始页，克隆页在其后保持顺序）。"""
    for _ in range(count):
        delete_slide(prs, 0)


# ---------------- 文字替换（保留样式） ----------------

def iter_text_shapes(shapes):
    """递归遍历（穿透 Group）产出所有带文字的形状。"""
    for shp in shapes:
        if shp.shape_type == 6:  # GROUP
            yield from iter_text_shapes(shp.shapes)
        elif getattr(shp, "has_text_frame", False):
            yield shp


def set_text_keep_style(text_frame, lines: list[str]) -> None:
    """替换整个文本框内容为 lines（每行一段），完整保留原第一段/第一 run 的字体样式。"""
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        return
    proto = copy.deepcopy(paragraphs[0]._p)  # 样式模板段

    # 清空既有段落，只留第一段
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)

    for i, line in enumerate(lines):
        if i == 0:
            p_el = paragraphs[0]._p
        else:
            p_el = copy.deepcopy(proto)
            text_frame._txBody.append(p_el)
        _set_paragraph_text(p_el, str(line))

    if not lines:
        _set_paragraph_text(paragraphs[0]._p, "")


def _set_paragraph_text(p_el, text: str) -> None:
    """段落内保留第一个 run（样式载体），仅改其文字；删除多余 run。"""
    runs = p_el.findall(qn("a:r"))
    if runs:
        for r in runs[1:]:
            p_el.remove(r)
        t = runs[0].find(qn("a:t"))
        if t is None:
            t = runs[0].makeelement(qn("a:t"), {})
            runs[0].append(t)
        t.text = text
    else:
        # 段落原本无 run（空段）：构造最简 run
        r = p_el.makeelement(qn("a:r"), {})
        t = r.makeelement(qn("a:t"), {})
        t.text = text
        r.append(t)
        # a:r 需要在 a:endParaRPr 之前
        end = p_el.find(qn("a:endParaRPr"))
        if end is not None:
            end.addprevious(r)
        else:
            p_el.append(r)


def biggest_text_shape(slide):
    """找到字号最大的文本形状（通常是标题）。"""
    best, best_size = None, -1.0
    for shp in iter_text_shapes(slide.shapes):
        size = _max_font_size(shp)
        if size > best_size and shp.text_frame.text.strip():
            best, best_size = shp, size
    return best


def _max_font_size(shp) -> float:
    size = 0.0
    for p in shp.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                size = max(size, r.font.size.pt)
    return size


_PAGE_NUM_RE = re.compile(r"^[\d\s/‹›<>]*$")
_MIN_TITLE_W = 2 * 914400  # 标题形状至少 2 英寸宽（排除页码等小占位符）


def find_title_shape(slide):
    """找内容框架页的标题形状：优先标题类占位符（TITLE=13/CENTER_TITLE=1/idx0），
    其次页面上部的宽文本形状。"""
    placeholder_hits = []
    for shp in slide.shapes:
        if getattr(shp, "is_placeholder", False) and shp.width and shp.width >= _MIN_TITLE_W:
            try:
                pf = shp.placeholder_format
                if pf.type in (13, 1, None) or pf.idx == 0:
                    placeholder_hits.append(shp)
            except Exception:
                pass
    if placeholder_hits:
        return max(placeholder_hits, key=lambda s: s.width)
    # 兜底：页面上部 1/4 区域内最宽的文本形状
    candidates = [s for s in iter_text_shapes(slide.shapes)
                  if s.top is not None and s.top < 1700000 and s.width and s.width >= _MIN_TITLE_W
                  and not _PAGE_NUM_RE.match(s.text_frame.text.strip() or "x")]
    return max(candidates, key=lambda s: s.width) if candidates else None


def shrink_text_to_shape(shp, min_size: int = 14, default_max: int = 26) -> None:
    """替换后的文字按形状宽度适配字号。

    - 有显式字号：从原字号逐级缩小到放得下（不放大，保模板风格）；
    - 继承字号（占位符未声明）：无法读取实际值，按形状宽高计算出
      [min_size, default_max] 区间内的合适字号并显式设置，
      避免细条标题带里出现过小/溢出文字。"""
    from pptx.util import Pt

    from app.ppt import text_engine
    if not shp.width:
        return
    box_w_pt = shp.width / 914400 * 72 - 8
    box_h_pt = (shp.height or 914400) / 914400 * 72
    for p in shp.text_frame.paragraphs:
        text = "".join(r.text for r in p.runs)
        if not text:
            continue
        sizes = [r.font.size.pt for r in p.runs if r.font.size]
        if sizes:
            size = max(sizes)
            while size > min_size and text_engine.text_width_pt(text, size) > box_w_pt:
                size -= 2
        else:
            size = min(default_max, max(min_size, int(box_h_pt * 0.55)))
            while size > min_size and text_engine.text_width_pt(text, size) > box_w_pt:
                size -= 2
        for r in p.runs:
            r.font.size = Pt(size)


_ENDING_TEXT_RE = re.compile(r"感谢|谢谢|thank|再见|the\s*end", re.I)


def clean_cloned_frame(slide, keep_shapes: list, zone_emu: tuple, page_area_emu: int) -> int:
    """清理克隆内容框架页上的模板残留示例内容（正文页绝不出现"谢谢"与示例图文）。

    删除规则：
    1. 任意位置含"感谢/谢谢/Thanks"类文本的形状（尾页内容泄漏到正文页的病根）；
    2. 正文区内的示例文本（>4 字符且与正文区重叠 >50%）；
    3. 正文区内的大图（占页面 >10% 且与正文区重叠 >30%——示例配图，保留必与新内容叠压）。
    保留：标题形状、页码占位符、页眉页脚小元素、边缘装饰。
    返回删除的形状数。
    """
    zx0, zy0, zx1, zy1 = zone_emu
    keep_ids = {id(s._element) for s in keep_shapes if s is not None}
    removed = 0

    def all_text(shp) -> str:
        if shp.shape_type == 6:
            return " ".join(s.text_frame.text for s in iter_text_shapes(shp.shapes))
        if getattr(shp, "has_text_frame", False):
            return shp.text_frame.text
        return ""

    def overlap_ratio(shp) -> float:
        try:
            if None in (shp.left, shp.top) or not shp.width or not shp.height:
                return 0.0
            ix = max(0, min(shp.left + shp.width, zx1) - max(shp.left, zx0))
            iy = max(0, min(shp.top + shp.height, zy1) - max(shp.top, zy0))
            return (ix * iy) / (shp.width * shp.height)
        except Exception:
            return 0.0

    for shp in list(slide.shapes):
        if id(shp._element) in keep_ids:
            continue
        try:
            if getattr(shp, "is_placeholder", False):
                try:
                    if shp.placeholder_format.type == 5:  # SLIDE_NUMBER 页码保留
                        continue
                except Exception:
                    pass
            text = all_text(shp).strip()
            ratio = overlap_ratio(shp)
            is_picture = shp.shape_type == 13 or (
                shp.shape_type == 6 and any(s.shape_type == 13 for s in shp.shapes))
            page_frac = ((shp.width or 0) * (shp.height or 0)) / max(page_area_emu, 1)

            if text and _ENDING_TEXT_RE.search(text):
                shp._element.getparent().remove(shp._element)
                removed += 1
            elif len(text.replace(" ", "")) > 2 and ratio > 0.3:
                # 覆盖"汇报人:"等 3~4 字残留（阈值从 >4字/50% 收紧到 >2字/30%）
                shp._element.getparent().remove(shp._element)
                removed += 1
            elif is_picture and page_frac > 0.10 and ratio > 0.3:
                shp._element.getparent().remove(shp._element)
                removed += 1
        except Exception:
            continue
    return removed


def clean_title_band(slide, title_shape, band_bottom_emu: int) -> int:
    """清理标题带内的模板残留文字（演示标题/"汇报人："等），只保留我们的标题。

    删除：完全位于正文区上方（top+height <= band_bottom）、含 >2 字符文本、
    且不是标题形状本身/页码占位符的文本形状。"""
    removed = 0
    title_id = id(title_shape._element) if title_shape is not None else None
    for shp in list(slide.shapes):
        try:
            if id(shp._element) == title_id:
                continue
            if getattr(shp, "is_placeholder", False):
                try:
                    if shp.placeholder_format.type == 5:
                        continue
                except Exception:
                    pass
            if shp.top is None or shp.height is None:
                continue
            if shp.top + shp.height > band_bottom_emu:
                continue
            text = " ".join(s.text_frame.text for s in iter_text_shapes([shp])).strip()
            if len(text.replace(" ", "")) > 2:
                shp._element.getparent().remove(shp._element)
                removed += 1
        except Exception:
            continue
    return removed


def remove_ending_texts(slide) -> int:
    """删除页面上的"感谢/谢谢"类文本形状（用于目录页等克隆页兜底）。"""
    removed = 0
    for shp in list(slide.shapes):
        try:
            texts = " ".join(s.text_frame.text for s in iter_text_shapes([shp]))
            if texts and _ENDING_TEXT_RE.search(texts):
                shp._element.getparent().remove(shp._element)
                removed += 1
        except Exception:
            continue
    return removed


_TOC_TITLE_WORDS = ("目录", "目 录", "contents", "content", "agenda")


def _is_toc_title(text: str) -> bool:
    return any(w in text.lower() for w in _TOC_TITLE_WORDS) and len(text) <= 10


def prepare_toc(slide, items: list[str]) -> tuple[str, tuple | None]:
    """目录页条目处理。返回 (mode, area)：
    - ("replaced", None)：顶层文本条目足够宽，已原位替换（样式全保留）；
    - ("cleared", (x,y,w,h) 英寸)：条目在 Group 内/过窄不可安全替换 →
      清除条目区形状（Group/线条/条目文本），返回该区域供自绘目录列表；
    - ("none", None)：未识别到条目结构，由调用方自选区域绘制。
    """
    from app.ppt import text_engine

    plain_entries, zone_shapes = [], []
    for shp in slide.shapes:  # 只看顶层：Group 整体视为条目构件
        st = shp.shape_type
        if st == 6:  # GROUP → 目录条目构件
            zone_shapes.append(shp)
        elif st in (9, 20):  # LINE / CONNECTOR 装饰线
            zone_shapes.append(shp)
        elif getattr(shp, "has_text_frame", False):
            text = shp.text_frame.text.strip()
            if not text or _is_toc_title(text) or _PAGE_NUM_RE.match(text):
                continue
            plain_entries.append(shp)

    # 路径 1：顶层纯文本条目原位替换（须数量够 + 12pt 单行可放下）
    if len(plain_entries) >= len(items) and not zone_shapes:
        plain_entries.sort(key=lambda s: (s.top or 0))
        fits = all(
            shp.width and text_engine.text_width_pt(item, 12) <= shp.width / 914400 * 72 - 6
            for shp, item in zip(plain_entries, items))
        if fits:
            for i, shp in enumerate(plain_entries):
                if i < len(items):
                    set_text_keep_style(shp.text_frame, [items[i]])
                    shrink_text_to_shape(shp, min_size=12)
                else:  # 多余条目整形删除（清空文字仍会残留装饰框）
                    shp._element.getparent().remove(shp._element)
            # 数字角标（01/02/03…）多于章节数时，按纵向顺序删除多余的
            chips = [s for s in iter_text_shapes(slide.shapes)
                     if re.fullmatch(r"0?\d{1,2}", (s.text_frame.text or "").strip())]
            chips.sort(key=lambda s: (s.top or 0))
            for extra in chips[len(items):]:
                try:
                    extra._element.getparent().remove(extra._element)
                except Exception:
                    pass
            logger.info("目录条目原位替换成功：%d 条（清理多余角标 %d 个）",
                        len(items), max(0, len(chips) - len(items)))
            return "replaced", None

    # 路径 2：彻底清除条目区后自绘（保留"目录"标题、页码与整页大装饰图）
    zone_shapes.extend(plain_entries)
    boxes = [(s.left, s.top, s.width, s.height) for s in zone_shapes
             if s.left is not None and s.top is not None and s.width and s.height]
    if not boxes:
        return "none", None
    zx0 = min(b[0] for b in boxes)
    zy0 = min(b[1] for b in boxes)
    zx1 = max(b[0] + b[2] for b in boxes)
    zy1 = max(b[1] + b[3] for b in boxes)

    # 彻底清理：与条目区重叠 >30% 的一切形状（含数字角标/色条/装饰块），
    # 避免残留与自绘目录叠压——只豁免目录标题、页码占位符、整页大图
    page_area = 12192000 * 6858000
    removed = 0
    for shp in list(slide.shapes):
        try:
            text = " ".join(s.text_frame.text for s in iter_text_shapes([shp])).strip()
            if text and _is_toc_title(text):
                continue
            if getattr(shp, "is_placeholder", False):
                try:
                    if shp.placeholder_format.type == 5:
                        continue
                except Exception:
                    pass
            if (shp.shape_type == 13 and shp.width and shp.height
                    and shp.width * shp.height > page_area * 0.2):
                continue  # 整页大装饰图保留
            if None in (shp.left, shp.top) or not shp.width or not shp.height:
                continue
            ix = max(0, min(shp.left + shp.width, zx1) - max(shp.left, zx0))
            iy = max(0, min(shp.top + shp.height, zy1) - max(shp.top, zy0))
            if (ix * iy) / (shp.width * shp.height) > 0.3:
                shp._element.getparent().remove(shp._element)
                removed += 1
        except Exception:
            continue
    logger.info("目录条目区彻底清理 %d 个形状，将在原区域自绘目录列表", removed)
    return "cleared", (zx0 / 914400, zy0 / 914400,
                      max((zx1 - zx0) / 914400, 4.5), max((zy1 - zy0) / 914400, 2.5))
