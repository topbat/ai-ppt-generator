"""受控版式绘制实现（确定性，无任何 LLM 参与）。

两层结构：
- BODY_PAINTERS：只绘制"正文区"的绘制器（x/y/w/h 由调用方给定）——
  模板克隆路径用它把内容画进模板内容框架页的正文区；
- LAYOUT_PAINTERS：整页绘制器（页头+页脚+正文），无可用模板页时的系统兜底。

美化原则：留白充分、层次分明（标题色>正文色>辅助色）、卡片化分块、
强调色来自模板主题、正文最小 12pt 可读性铁律。
"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.core.logging import get_logger
from app.ppt import chart_builder, table_builder, text_engine
from app.ppt import design_tokens as dt
from app.ppt.fonts import set_run_font

logger = get_logger(__name__)

# ---- 页面几何常量（英寸，系统整页路径用；取值来自 Design Token） ----
PAGE_W, PAGE_H = dt.CANVAS_W, dt.CANVAS_H
MARGIN = dt.MARGIN_X + 0.2   # 系统整页路径左右边距（锚线 0.7 + 视觉缓冲）
CONTENT_X = MARGIN
CONTENT_W = PAGE_W - MARGIN * 2
CONTENT_BOTTOM = PAGE_H - dt.MARGIN_BOTTOM - 0.1

BODY_SIZE = dt.BODY_SIZE_BY_DENSITY   # 正文 16/18/20（演示媒介，拒绝 Word 化小字）
MIN_BODY_SIZE = dt.MIN_BODY_SIZE
GRAY_TEXT = dt.MUTED_COLOR            # 辅助文字（非纯黑灰阶体系）
LIGHT_BORDER = dt.BORDER_COLOR
CARD_BG = dt.SURFACE_COLOR


def _tint(hex_color: str, factor: float) -> str:
    r, g, b = (int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
    mix = lambda c: int(c + (255 - c) * factor)  # noqa: E731
    return f"{mix(r):02X}{mix(g):02X}{mix(b):02X}"


def _rect(slide, x, y, w, h, fill_hex, line_hex=None, rounded=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if line_hex:
        shape.line.color.rgb = RGBColor.from_string(line_hex)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if rounded:  # 收小圆角半径，更利落
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def _text(rc, slide, text, x, y, w, h, size, color, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None, min_size=MIN_BODY_SIZE,
          line_spacing=1.28, page=None, role=None):
    text = str(text or "")
    font = font or rc.theme.font_body
    if role in dt.TYPO_ROLES:
        min_size = max(min_size, dt.TYPO_ROLES[role][0])
    fitted, overflow = text_engine.fit_size(text, w * 72, h * 72, size, min_size, line_spacing)
    if overflow:
        text = text_engine.truncate_to_fit(text, w * 72, h * 72, min_size, line_spacing)
        rc.note(page, "text_truncated", f"文本超出容器被截断: {text[:20]}…")
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_run_font(run, font, fitted, color, bold)
    return tb


def _bullets(rc, slide, items, x, y, w, h, size, color, page=None, accent=None):
    """要点列表：强调色符号 + 正文双 run，自动降字号。

    最小字号仍放不下时确定性精简（先减条目保底 2 条，再截短文本），
    保证绝不越界溢出；精简行为记 info 级留痕。"""
    items = [str(i) for i in (items or []) if str(i).strip()]
    if not items:
        return
    orig_count = len(items)

    def measure(lst):
        return text_engine.fit_size("\n".join(lst), (w - 0.35) * 72,
                                    (h - 0.12 * len(lst)) * 72, size, MIN_BODY_SIZE, 1.45)

    fitted, overflow = measure(items)
    while overflow and len(items) > 2:      # 1) 减条目
        items = items[:-1]
        fitted, overflow = measure(items)
    guard = 0
    while overflow and guard < 6:           # 2) 截短文本
        items = [i if len(i) <= 10 else i[: max(10, int(len(i) * 0.78))] + "…" for i in items]
        fitted, overflow = measure(items)
        guard += 1
    if len(items) < orig_count or guard:
        rc.note(page, "bullets_trimmed",
                f"要点精简以适配版面（{orig_count}→{len(items)}条）")
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.4
        p.space_after = Pt(10)
        glyph = p.add_run()
        glyph.text = "▪  "
        set_run_font(glyph, rc.theme.font_body, fitted, accent or rc.accent_color, bold=True)
        run = p.add_run()
        run.text = item
        set_run_font(run, rc.theme.font_body, fitted, color)


def _header(rc, slide, sd) -> float:
    """系统路径页头：主色竖条 + 标题 + 细分隔线 + 页码。返回正文区起始 Y。"""
    theme = rc.theme
    _rect(slide, MARGIN, 0.52, 0.10, 0.52, theme.primary)
    _text(rc, slide, sd.get("title", ""), MARGIN + 0.28, 0.40, CONTENT_W - 0.28, 0.72,
          dt.TYPO["title"]["size"], theme.secondary, bold=True, font=theme.font_title,
          min_size=dt.TYPO["title"]["min"], page=sd.get("page"))
    _rect(slide, MARGIN, 1.22, CONTENT_W, 0.014, LIGHT_BORDER)
    y = 1.45
    if sd.get("subtitle"):
        _text(rc, slide, sd["subtitle"], MARGIN + 0.28, 1.32, CONTENT_W - 0.28, 0.45,
              13, GRAY_TEXT, page=sd.get("page"))
        y = 1.90
    _text(rc, slide, str(sd.get("page", "")), PAGE_W - 1.1, 7.08, 0.7, 0.32,
          10, GRAY_TEXT, align=PP_ALIGN.RIGHT)
    return y


# ================= 元素绘制器（区域参数化） =================

def _el_bullet_group(rc, slide, el, x, y, w, h, page):
    _bullets(rc, slide, el.get("items", []), x, y + 0.1, w, h - 0.1,
             BODY_SIZE[rc.density], rc.body_color, page)


def _el_cards(rc, slide, el, x, y, w, h, page):
    items = (el.get("items") or [])[:6]
    if not items:
        return
    n = len(items)
    cols = n if n <= 3 else (2 if n == 4 else 3)
    rows = (n + cols - 1) // cols
    gap = 0.3
    cw, chh = (w - gap * (cols - 1)) / cols, (h - gap * (rows - 1)) / rows
    # 卡片高度按内容自适应（单行时避免整版拉伸导致大面积留白），并整体垂直居中
    if rows == 1:
        max_desc = max((len(str(i.get("desc", ""))) for i in items), default=0)
        chh = min(chh, 1.7 + 0.055 * max_desc)
        y = y + max(0, (h - chh) * 0.38)
    for i, item in enumerate(items):
        cx = x + (i % cols) * (cw + gap)
        cy = y + (i // cols) * (chh + gap)
        _rect(slide, cx, cy, cw, chh, CARD_BG, line_hex=LIGHT_BORDER, rounded=True)
        _rect(slide, cx + 0.02, cy, cw - 0.04, 0.06, rc.theme.primary)
        _text(rc, slide, item.get("title", ""), cx + 0.25, cy + 0.22, cw - 0.5, 0.5,
              16, rc.theme.secondary, bold=True, min_size=13, page=page,
              font=rc.theme.font_title)
        _text(rc, slide, item.get("desc", ""), cx + 0.25, cy + 0.8, cw - 0.5, chh - 1.05,
              12.5, GRAY_TEXT, page=page, line_spacing=1.4)


def _el_key_number(rc, slide, el, x, y, w, h, page):
    items = (el.get("items") or [])[:4]
    if not items:
        return
    cw = w / len(items)
    for i, item in enumerate(items):
        cx = x + i * cw
        if i > 0:  # 项间细分隔线
            _rect(slide, cx, y + h * 0.18, 0.014, h * 0.55, rc.divider_color)
        _text(rc, slide, item.get("value", ""), cx, y + h * 0.10, cw, h * 0.45,
              dt.TYPO["number"]["size"], rc.accent_color, bold=True, align=PP_ALIGN.CENTER,
              min_size=dt.TYPO["number"]["min"], page=page, font=rc.theme.font_title)
        _text(rc, slide, item.get("label", ""), cx + 0.15, y + h * 0.60, cw - 0.3, h * 0.3,
              14, rc.muted_color, align=PP_ALIGN.CENTER, page=page)


def _el_chart(rc, slide, el, x, y, w, h, page):
    try:
        # 深底页图表坐标轴/标签用浅色，保证可读
        chart_builder.add_chart(slide, el, Inches(x), Inches(y), Inches(w), Inches(h), rc.theme,
                                font_color="F2F5FA" if rc.on_dark else None,
                                series_color=rc.accent_color)
    except Exception as e:
        logger.warning("第 %s 页图表构建失败，降级为表格: %s", page, e)
        rc.note(page, "chart_degraded", "图表数据异常，已降级为表格")
        data = chart_builder.sanitize_chart_data(el.get("data", []))
        if data:
            table_builder.add_table(slide, ["项目", f"数值({el.get('unit', '')})"],
                                    [[d["label"], d["value"]] for d in data],
                                    Inches(x), Inches(y), Inches(w), Inches(h), rc.theme)
        else:
            _bullets(rc, slide, [el.get("title", "数据暂缺")], x, y, w, h, 15, rc.body_color, page)


def _el_table(rc, slide, el, x, y, w, h, page):
    try:
        table_builder.add_table(slide, el.get("headers", []), el.get("rows", []),
                                Inches(x), Inches(y), Inches(w), Inches(h), rc.theme)
    except Exception as e:
        logger.warning("第 %s 页表格构建失败，降级为要点: %s", page, e)
        rc.note(page, "table_degraded", "表格数据异常，已降级为要点")
        rows = ["  ".join(str(c) for c in r) for r in (el.get("rows") or [])[:6] if isinstance(r, list)]
        _bullets(rc, slide, rows, x, y, w, h, 13, rc.theme.text, page)


def _el_timeline(rc, slide, el, x, y, w, h, page):
    items = (el.get("items") or [])[:6]
    if not items:
        return
    line_y = y + h * 0.42
    _rect(slide, x + 0.2, line_y, w - 0.4, 0.03, _tint(rc.theme.primary, 0.55))
    seg = (w - 0.4) / len(items)
    # 列宽不足以承载描述时自动省略（时间+标题更干净，避免截断与叠压）
    show_desc = seg >= 2.0 and len(items) <= 3
    if not show_desc and any(it.get("desc") for it in items if isinstance(it, dict)):
        rc.note(page, "timeline_desc_dropped", "时间轴节点较密，已省略节点描述保持版面整洁")
    for i, item in enumerate(items):
        cx = x + 0.2 + seg * i + seg / 2
        outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.11), Inches(line_y - 0.095),
                                       Inches(0.22), Inches(0.22))
        outer.fill.solid()
        outer.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        outer.line.color.rgb = RGBColor.from_string(rc.accent_color)
        outer.line.width = Pt(1.75)
        outer.shadow.inherit = False
        _text(rc, slide, item.get("time", ""), cx - seg / 2 + 0.1, line_y - 0.72, seg - 0.2, 0.42,
              15, rc.accent_color, bold=True, align=PP_ALIGN.CENTER, min_size=11, page=page,
              font=rc.theme.font_title)
        _text(rc, slide, item.get("title", ""), cx - seg / 2 + 0.1, line_y + 0.28, seg - 0.2, 0.75,
              13.5, rc.heading_color, bold=True, align=PP_ALIGN.CENTER, min_size=11, page=page)
        if show_desc and item.get("desc"):
            _text(rc, slide, item["desc"], cx - seg / 2 + 0.15, line_y + 1.05, seg - 0.3,
                  max(0.6, h * 0.3), 11.5, rc.muted_color, align=PP_ALIGN.CENTER,
                  min_size=10, page=page)


def _el_process(rc, slide, el, x, y, w, h, page):
    steps = [str(s) for s in (el.get("steps") or [])[:6]]
    if not steps:
        return
    gap = 0.14
    cw = (w - gap * (len(steps) - 1)) / len(steps)
    cy, chh = y + h * 0.26, min(h * 0.48, 1.35)
    for i, step in enumerate(steps):
        cx = x + i * (cw + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx), Inches(cy), Inches(cw), Inches(chh))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(_tint(rc.theme.primary, 0.10 * i))
        shape.line.fill.background()
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step
        size, _ = text_engine.fit_size(step, (cw - 0.35) * 72, chh * 72, 14, 10)
        set_run_font(run, rc.theme.font_body, size, "FFFFFF", bold=True)


def _el_architecture(rc, slide, el, x, y, w, h, page):
    layers = (el.get("layers") or [])[:5]
    if not layers:
        return
    gap = 0.16
    lh = (h - gap * (len(layers) - 1)) / len(layers)
    for i, layer in enumerate(layers):
        ly = y + i * (lh + gap)
        fill = _tint(rc.theme.primary, 0.08 + 0.15 * i)
        _rect(slide, x, ly, w, lh, fill, rounded=True)
        _text(rc, slide, layer.get("title", ""), x + 0.3, ly, 1.9, lh,
              15, "FFFFFF", bold=True, min_size=12, anchor=MSO_ANCHOR.MIDDLE, page=page,
              font=rc.theme.font_title)
        items = (layer.get("items") or [])[:5]
        if items:
            chip_x = x + 2.4
            chip_w = (w - 2.75 - 0.18 * (len(items) - 1)) / len(items)
            for j, item in enumerate(items):
                cx = chip_x + j * (chip_w + 0.18)
                chip = _rect(slide, cx, ly + lh * 0.18, chip_w, lh * 0.64, "FFFFFF", rounded=True)
                tf = chip.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                size, _ = text_engine.fit_size(str(item), (chip_w - 0.12) * 72, lh * 0.64 * 72, 12, 9)
                run.text = str(item)
                set_run_font(run, rc.theme.font_body, size, rc.theme.secondary)


def _el_comparison(rc, slide, el, x, y, w, h, page):
    half = (w - 0.35) / 2
    for i, key in enumerate(("left", "right")):
        side = el.get(key) or {}
        sx = x + i * (half + 0.35)
        head_color = "8C8C8C" if i == 0 else rc.theme.primary
        _rect(slide, sx, y, half, h, CARD_BG, line_hex=LIGHT_BORDER, rounded=True)
        _rect(slide, sx, y, half, 0.58, head_color)
        _text(rc, slide, side.get("title", ""), sx, y + 0.08, half, 0.45,
              16, "FFFFFF", bold=True, align=PP_ALIGN.CENTER, min_size=12, page=page,
              font=rc.theme.font_title)
        _bullets(rc, slide, side.get("items", []), sx + 0.3, y + 0.82, half - 0.6, h - 1.05,
                 BODY_SIZE[rc.density] - 1, rc.theme.text, page,
                 accent="BFBFBF" if i == 0 else None)


def _el_quote(rc, slide, el, x, y, w, h, page):
    _text(rc, slide, "“", x, y - 0.35, 1.4, 1.4, 80, _tint(rc.theme.primary, 0.45),
          bold=True, font=rc.theme.font_title)
    _text(rc, slide, el.get("text", ""), x + 0.9, y + 0.7, w - 1.8, h - 1.5,
          22, rc.heading_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          min_size=14, page=page, font=rc.theme.font_title, line_spacing=1.5)
    if el.get("author"):
        _text(rc, slide, f"—— {el['author']}", x, y + h - 0.55, w - 0.6, 0.5,
              13, rc.muted_color, align=PP_ALIGN.RIGHT, page=page)


def _el_paragraph(rc, slide, el, x, y, w, h, page):
    _text(rc, slide, el.get("text", ""), x, y + 0.1, w, h - 0.1, BODY_SIZE[rc.density],
          rc.body_color, line_spacing=1.6, page=page)


ELEMENT_PAINTERS = {
    "bullet_group": _el_bullet_group, "cards": _el_cards, "key_number": _el_key_number,
    "chart": _el_chart, "table": _el_table, "timeline": _el_timeline, "process": _el_process,
    "architecture": _el_architecture, "comparison": _el_comparison, "quote": _el_quote,
    "paragraph": _el_paragraph,
}


def _paint_stacked(rc, slide, sd, elements, x, y, w, h):
    elements = [el for el in elements if el.get("type") in ELEMENT_PAINTERS][:3]
    if not elements:
        return
    weights = [3 if el["type"] in ("chart", "table", "architecture", "timeline", "comparison") else
               max(1, len(el.get("items", el.get("steps", [1, 1])))) for el in elements]
    total = sum(weights)
    cy = y
    for el, wt in zip(elements, weights):
        eh = h * wt / total
        ELEMENT_PAINTERS[el["type"]](rc, slide, el, x, cy, w, eh - 0.15, sd.get("page"))
        cy += eh


def _first_el(sd, *types):
    for el in sd.get("elements", []):
        if el.get("type") in types:
            return el
    return None


def _coerce_cards(sd, count):
    el = _first_el(sd, "cards")
    if el and el.get("items"):
        return {"type": "cards", "items": el["items"][:count]}
    bullets = _first_el(sd, "bullet_group")
    if bullets and bullets.get("items"):
        items = bullets["items"][:count]
        return {"type": "cards", "items": [
            {"title": it[:12], "desc": it[12:] if len(it) > 12 else ""} for it in items]}
    return None


# ================= 正文区绘制器（BODY_PAINTERS，模板克隆路径复用） =================

def body_title_content(rc, slide, sd, x, y, w, h):
    _paint_stacked(rc, slide, sd, sd.get("elements", []), x, y, w, h)


def body_two_column(rc, slide, sd, x, y, w, h):
    half = (w - 0.55) / 2
    els = [el for el in sd.get("elements", []) if el.get("type") in ELEMENT_PAINTERS]
    if len(els) >= 2:
        _paint_stacked(rc, slide, sd, [els[0]], x, y, half, h)
        _paint_stacked(rc, slide, sd, els[1:2], x + half + 0.55, y, half, h)
    elif len(els) == 1 and els[0]["type"] == "bullet_group":
        items = els[0].get("items", [])
        mid = (len(items) + 1) // 2
        _el_bullet_group(rc, slide, {"items": items[:mid]}, x, y, half, h, sd.get("page"))
        _el_bullet_group(rc, slide, {"items": items[mid:]}, x + half + 0.55, y, half, h, sd.get("page"))
    else:
        _paint_stacked(rc, slide, sd, els, x, y, w, h)


def _body_cards(count):
    def painter(rc, slide, sd, x, y, w, h):
        el = _coerce_cards(sd, count)
        if el:
            _el_cards(rc, slide, el, x, y + 0.05, w, h - 0.05, sd.get("page"))
        else:
            _paint_stacked(rc, slide, sd, sd.get("elements", []), x, y, w, h)
    return painter


def _body_single_el(el_type):
    def painter(rc, slide, sd, x, y, w, h):
        el = _first_el(sd, el_type)
        if el:
            ELEMENT_PAINTERS[el_type](rc, slide, el, x, y + 0.1, w, h - 0.2, sd.get("page"))
        else:
            rc.note(sd.get("page"), "element_missing", f"缺少 {el_type} 元素，已按通用版式绘制")
            _paint_stacked(rc, slide, sd, sd.get("elements", []), x, y, w, h)
    return painter


def _body_chart(chart_type):
    def painter(rc, slide, sd, x, y, w, h):
        el = _first_el(sd, "chart")
        if el:
            el = {**el, "chart_type": el.get("chart_type") or chart_type}
            _el_chart(rc, slide, el, x + 0.7, y + 0.1, w - 1.4, h - 0.25, sd.get("page"))
        else:
            rc.note(sd.get("page"), "element_missing", "缺少 chart 元素，已按通用版式绘制")
            _paint_stacked(rc, slide, sd, sd.get("elements", []), x, y, w, h)
    return painter


def body_key_number(rc, slide, sd, x, y, w, h):
    el = _first_el(sd, "key_number")
    if el:
        _el_key_number(rc, slide, el, x, y + h * 0.15, w, h * 0.66, sd.get("page"))
    else:
        _paint_stacked(rc, slide, sd, sd.get("elements", []), x, y, w, h)


def body_quote(rc, slide, sd, x, y, w, h):
    el = _first_el(sd, "quote") or {"text": sd.get("key_message") or sd.get("title", "")}
    _el_quote(rc, slide, el, x + 0.9, y + 0.15, w - 1.8, h - 0.3, sd.get("page"))


def body_section(rc, slide, sd, x, y, w, h):
    """章节页正文：垂直居中的固定构图块（序号→分隔线→标题→核心观点），
    构图与间距固定，保证同一份 PPT 内所有章节页、不同模板之间位置一致。"""
    no = sd.get("section_no") or 1
    theme = rc.theme
    block_h = 4.0 if sd.get("key_message") else 3.4
    top = y + max(0, (h - block_h) * 0.44)  # 居中略偏上，符合视觉重心
    _text(rc, slide, f"{no:02d}", x + 0.3, top, 3.2, 1.7, 88,
          _tint(theme.primary, 0.6 if rc.on_dark else 0.25), bold=True,
          font=theme.font_title, min_size=48)
    _rect(slide, x + 0.42, top + 1.9, 1.5, 0.05, rc.accent_color)
    _text(rc, slide, sd.get("title", ""), x + 0.42, top + 2.12, w - 1.0, 0.85,
          30, rc.heading_color, bold=True, font=theme.font_title, min_size=20, page=sd.get("page"))
    if sd.get("key_message"):
        km = str(sd["key_message"])
        if len(km) > 48:  # 核心观点超长时确定性精简（info 留痕，不产生截断告警）
            km = km[:47] + "…"
            rc.note(sd.get("page"), "key_message_trimmed", "章节核心观点超长已精简")
        _text(rc, slide, km, x + 0.45, top + 3.1, w - 1.2, 0.85,
              13.5, rc.muted_color, page=sd.get("page"), line_spacing=1.35)


def body_summary(rc, slide, sd, x, y, w, h):
    el = _first_el(sd, "bullet_group")
    if not el:
        return
    items = el.get("items", [])
    # 垂直居中 + 左侧留白，避免整页上浮空旷
    needed = 0.62 * len(items)
    cy = y + max(0, (h - needed) * 0.4)
    _bullets(rc, slide, items, x + 0.5, cy, w - 1.0, needed + 0.4,
             17, rc.body_color, sd.get("page"))


def body_toc(rc, slide, sd, x, y, w, h):
    """目录正文（模板目录条目替换失败时的兜底绘制）。"""
    el = _first_el(sd, "bullet_group") or {"items": []}
    items = el.get("items", [])[:10]
    cols = 2 if len(items) > 6 else 1
    per = (len(items) + cols - 1) // cols if items else 0
    cw = (w - 0.8) / cols if cols else w
    for c in range(cols):
        seg = items[c * per:(c + 1) * per]
        cy = y + 0.25
        for i, item in enumerate(seg):
            idx = c * per + i + 1
            _text(rc, slide, f"{idx:02d}", x + c * (cw + 0.8), cy, 0.85, 0.6,
                  22, rc.accent_color, bold=True, font=rc.theme.font_title)
            _text(rc, slide, str(item), x + c * (cw + 0.8) + 0.95, cy + 0.07, cw - 1.05, 0.6,
                  17, rc.heading_color, page=sd.get("page"))
            _rect(slide, x + c * (cw + 0.8) + 0.95, cy + 0.62, cw - 1.5, 0.012, rc.divider_color)
            cy += 0.85


BODY_PAINTERS = {
    "title_content": body_title_content, "two_column": body_two_column,
    "three_column": _body_cards(3), "three_cards": _body_cards(3),
    "four_cards": _body_cards(4), "key_number": body_key_number, "quote": body_quote,
    "timeline": _body_single_el("timeline"), "process": _body_single_el("process"),
    "architecture": _body_single_el("architecture"),
    "comparison": _body_single_el("comparison"), "table": _body_single_el("table"),
    "bar_chart": _body_chart("bar"), "line_chart": _body_chart("line"),
    "pie_chart": _body_chart("pie"), "image_text": body_title_content,
    "section": body_section, "summary": body_summary, "toc": body_toc,
}


# ================= 整页绘制器（系统兜底路径） =================

def paint_cover(rc, slide, sd):
    theme = rc.theme
    _rect(slide, 0, 0, PAGE_W, PAGE_H, theme.primary)
    _rect(slide, 0, PAGE_H - 0.9, PAGE_W, 0.05, _tint(theme.primary, 0.45))
    _rect(slide, 1.25, 2.35, 1.6, 0.06, _tint(theme.primary, 0.6))
    _text(rc, slide, sd.get("title", ""), 1.2, 2.7, PAGE_W - 2.4, 1.6, 40, "FFFFFF",
          bold=True, font=theme.font_title, min_size=26, page=sd.get("page"))
    if sd.get("subtitle"):
        _text(rc, slide, sd["subtitle"], 1.25, 4.35, PAGE_W - 2.5, 0.8, 17,
              _tint(theme.primary, 0.7), page=sd.get("page"))


def paint_toc(rc, slide, sd):
    y = _header(rc, slide, {**sd, "title": sd.get("title") or "目录"})
    body_toc(rc, slide, sd, CONTENT_X, y + 0.15, CONTENT_W, CONTENT_BOTTOM - y - 0.15)


def paint_section(rc, slide, sd):
    theme = rc.theme
    _rect(slide, 0, 0, PAGE_W, PAGE_H, theme.primary)
    no = sd.get("section_no") or 1
    _text(rc, slide, f"{no:02d}", 1.2, 1.9, 3.2, 1.9, 72, _tint(theme.primary, 0.5),
          bold=True, font=theme.font_title)
    _rect(slide, 1.3, 3.95, 1.6, 0.05, _tint(theme.primary, 0.55))
    _text(rc, slide, sd.get("title", ""), 1.25, 4.2, PAGE_W - 2.5, 1.0, 32, "FFFFFF",
          bold=True, font=theme.font_title, min_size=22, page=sd.get("page"))
    if sd.get("key_message"):
        _text(rc, slide, sd["key_message"], 1.3, 5.35, PAGE_W - 2.6, 0.7, 15,
              _tint(theme.primary, 0.72), page=sd.get("page"))


def paint_summary(rc, slide, sd):
    theme = rc.theme
    _rect(slide, 0, 0, PAGE_W, PAGE_H, theme.primary)
    _text(rc, slide, sd.get("title") or "总结", 1.25, 0.9, PAGE_W - 2.5, 1.0, 32, "FFFFFF",
          bold=True, font=theme.font_title, page=sd.get("page"))
    _rect(slide, 1.3, 2.0, 1.5, 0.05, _tint(theme.primary, 0.55))
    el = _first_el(sd, "bullet_group")
    if el:
        _bullets(rc, slide, el.get("items", []), 1.3, 2.45, PAGE_W - 2.6, 4.1,
                 17, "FFFFFF", sd.get("page"), accent=_tint(theme.primary, 0.6))


def paint_ending(rc, slide, sd):
    theme = rc.theme
    _rect(slide, 0, 0, PAGE_W, PAGE_H, theme.primary)
    _text(rc, slide, sd.get("title") or "感谢聆听", 1.2, 2.9, PAGE_W - 2.4, 1.4, 48, "FFFFFF",
          bold=True, align=PP_ALIGN.CENTER, font=theme.font_title, page=sd.get("page"))


def _full_page(body_key):
    def painter(rc, slide, sd):
        y = _header(rc, slide, sd)
        BODY_PAINTERS[body_key](rc, slide, sd, CONTENT_X, y + 0.15,
                                CONTENT_W, CONTENT_BOTTOM - y - 0.15)
    return painter


LAYOUT_PAINTERS = {
    "cover": paint_cover, "toc": paint_toc, "section": paint_section,
    "summary": paint_summary, "ending": paint_ending,
    **{k: _full_page(k) for k in ("title_content", "two_column", "three_column", "three_cards",
                                  "four_cards", "key_number", "quote", "timeline", "process",
                                  "architecture", "comparison", "table", "bar_chart",
                                  "line_chart", "pie_chart", "image_text")},
}
paint_title_content = LAYOUT_PAINTERS["title_content"]  # renderer 降级兜底引用

