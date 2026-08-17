"""PPT 渲染总控：Presentation JSON → PPTX（确定性，无 LLM）。

双路径渲染：
1. 模板克隆路径（首选）：模板解析出可用角色页（cover/toc/content_frame/ending）时，
   直接克隆模板页保留全部视觉（背景/装饰/图片/字体），只替换文字/绘制正文；
2. 系统路径（兜底）：模板无可用页时，用主题 Token + 系统版式整页绘制。

兜底铁律（文档 §66-8）：单页绘制失败绝不导致整册失败。
"""
import time

from pptx import Presentation
from pptx.util import Emu, Inches

from app.core.errors import PPTError
from app.core.logging import get_logger
from app.ppt import template_clone as tc
from app.ppt.layouts import (BODY_PAINTERS, GRAY_TEXT, LAYOUT_PAINTERS, PAGE_H, PAGE_W, _text,
                             paint_title_content)
from app.schemas.presentation import PresentationSpec, ThemeSpec

logger = get_logger(__name__)

_EMU_IN = 914400.0


class RenderContext:
    """渲染上下文。on_dark 为当前页深色背景标记（来自模板背景识别），
    所有正文配色通过下列属性取值，保证"深底浅字、浅底深字"的协调与醒目。"""

    def __init__(self, theme: ThemeSpec, density: str):
        self.theme = theme
        self.density = density
        self.on_dark = False
        self.notes: list[dict] = []

    def note(self, page, issue: str, detail: str):
        self.notes.append({"page": page, "issue": issue, "detail": detail})

    # ---- 自适应配色（取值来自 Design Token） ----
    @property
    def heading_color(self) -> str:
        return "FFFFFF" if self.on_dark else self.theme.secondary

    @property
    def body_color(self) -> str:
        from app.ppt import design_tokens as dt
        return dt.DARK_TEXT_COLOR if self.on_dark else dt.TEXT_COLOR

    @property
    def muted_color(self) -> str:
        from app.ppt import design_tokens as dt
        return dt.DARK_MUTED_COLOR if self.on_dark else dt.MUTED_COLOR

    @property
    def accent_color(self) -> str:
        from app.ppt.layouts import _tint
        # 深底上主题色可能不够亮，提亮后再用，保证醒目
        return _tint(self.theme.primary, 0.45) if self.on_dark else self.theme.primary

    @property
    def divider_color(self) -> str:
        from app.ppt import design_tokens as dt
        return dt.DARK_BORDER_COLOR if self.on_dark else dt.BORDER_COLOR


def render_presentation(spec: PresentationSpec, out_path: str,
                        template: dict | None = None) -> dict:
    """template: {"path": str, "slides": [{"slide_index", "role", "confidence", "title_geo"}]}
    返回 {"notes", "degraded_pages", "duration_ms", "render_path"}"""
    started = time.monotonic()
    rc = RenderContext(spec.theme, spec.density)

    roles = _pick_role_slides(template) if template else {}
    use_template = bool(roles.get("content_frame") is not None or roles.get("cover") is not None)

    if use_template:
        degraded = _render_with_template(rc, spec, template["path"], roles, out_path)
        render_path = "template_clone"
    else:
        degraded = _render_with_system(rc, spec, out_path)
        render_path = "system"

    duration_ms = int((time.monotonic() - started) * 1000)
    logger.info("渲染完成[%s]：%d 页 / 耗时 %dms / 降级页 %s / 备注 %d 条",
                render_path, len(spec.slides), duration_ms, degraded or "无", len(rc.notes))
    return {"notes": rc.notes, "degraded_pages": degraded,
            "duration_ms": duration_ms, "render_path": render_path}


def _pick_role_slides(template: dict | None) -> dict:
    """每个角色取置信度最高的模板页。"""
    roles: dict[str, dict | None] = {}
    for s in (template or {}).get("slides", []):
        role = s.get("role")
        if role in ("cover", "toc", "content_frame", "ending", "section"):
            cur = roles.get(role)
            if cur is None or s.get("confidence", 0) > cur.get("confidence", 0):
                roles[role] = s
    return roles


# ================= 模板克隆路径 =================

def _render_with_template(rc, spec, tpl_path: str, roles: dict, out_path: str) -> list[int]:
    degraded: list[int] = []
    try:
        prs = Presentation(tpl_path)
    except Exception as e:
        logger.error("模板打开失败，回退系统路径: %s", e)
        return _render_with_system(rc, spec, out_path)

    orig_count = len(prs.slides._sldIdLst)
    page_w = Emu(prs.slide_width).inches if prs.slide_width else PAGE_W
    page_h = Emu(prs.slide_height).inches if prs.slide_height else PAGE_H
    slides_by_index = {s.get("slide_index"): s for s in roles.values() if s}
    blank_layout = _emptiest_layout(prs)

    def src_slide(role: str):
        info = roles.get(role)
        return prs.slides[info["slide_index"]] if info else None

    def content_area(role_info) -> tuple[float, float, float, float]:
        """(x, y, w, h) 正文绘制区：标题形状下方 → 页底留白之上。

        几何合法性兜底：模板标题占位符位置异常（过低/过窄）会把正文区算得
        极小导致整页截断，此时退回默认安全区域。"""
        contract = (role_info or {}).get("space_contract") or {}
        safe = contract.get("safe_zone")
        if safe:
            try:
                return (float(safe["x"]), float(safe["y"]),
                        float(safe["width"]), float(safe["height"]))
            except (KeyError, TypeError, ValueError):
                logger.warning("模板空间契约无效，退回标题几何计算")
        geo = (role_info or {}).get("title_geo")
        if geo and geo.get("width"):
            x = max(0.55, geo["left"] / _EMU_IN)
            y = (geo["top"] + geo["height"]) / _EMU_IN + 0.28
            w = min(geo["width"] / _EMU_IN, page_w - 2 * x + 0.4)
        else:
            x, y, w = 0.9, 1.5, page_w - 1.8
        h = page_h - 0.55 - y
        if y > page_h * 0.45 or h < 2.2 or w < 5.0:
            x, y, w = 0.9, min(1.5, page_h * 0.2), page_w - 1.8
            h = page_h - 0.55 - y
        return x, y, w, h

    logger.info("模板克隆渲染：角色页=%s", {k: v["slide_index"] for k, v in roles.items() if v})
    for sd in sorted(spec.slides, key=lambda s: s.page):
        data = sd.model_dump()
        rc.on_dark = False  # 每页重置，具体分支按来源模板页背景设置
        try:
            _render_template_page(rc, prs, sd, data, roles, src_slide, content_area, blank_layout)
        except Exception as e:
            logger.warning("第 %d 页模板克隆失败，降级为系统版式: %s", sd.page, e)
            degraded.append(sd.page)
            rc.note(sd.page, "page_degraded", f"模板克隆失败已降级: {e}")
            try:
                slide = prs.slides.add_slide(blank_layout)
                LAYOUT_PAINTERS.get(sd.type, paint_title_content)(rc, slide, data)
            except Exception as e2:
                logger.error("第 %d 页降级绘制仍失败(保留空页): %s", sd.page, e2)
                prs.slides.add_slide(blank_layout)

    tc.delete_leading_slides(prs, orig_count)  # 删除模板原始页，克隆页保持顺序

    # 终检：非尾页绝不允许出现"感谢/谢谢"类内容（质量标准 content_sanity 维度）
    # 页面层与版式层一并扫描（版式层艺术字同样会被渲染出来）
    ordered = sorted(spec.slides, key=lambda s: s.page)
    for idx, slide in enumerate(prs.slides):
        sd = ordered[idx] if idx < len(ordered) else None
        if sd is None or sd.type == "ending":
            continue
        texts = " ".join(s.text_frame.text for s in tc.iter_text_shapes(slide.shapes))
        try:
            texts += " " + " ".join(s.text_frame.text
                                    for s in tc.iter_text_shapes(slide.slide_layout.shapes))
        except Exception:
            pass
        if texts.strip() and tc._ENDING_TEXT_RE.search(texts):
            rc.note(sd.page, "ending_leak", "非尾页出现感谢/谢谢类尾页内容")

    _save(prs, out_path)
    return degraded


def _render_template_page(rc, prs, sd, data, roles, src_slide, content_area, blank_layout):
    stype = sd.type
    frame_info = roles.get("content_frame")

    if stype == "cover" and src_slide("cover") is not None:
        slide = tc.clone_slide(prs, src_slide("cover"))
        title_shape = tc.biggest_text_shape(slide)
        page_w = Emu(prs.slide_width).inches if prs.slide_width else PAGE_W
        page_h = Emu(prs.slide_height).inches if prs.slide_height else PAGE_H
        # 模板封面标题为文本形状且字号足够大 → 原位替换（风格全保留）
        if title_shape is not None and tc._max_font_size(title_shape) >= 18:
            lines = [sd.title]
            if sd.subtitle and len(sd.title) <= 14:
                lines.append(sd.subtitle)
            tc.set_text_keep_style(title_shape.text_frame, lines)
            # 副标题行缩小弱化（不与主标题同级）
            if len(lines) > 1:
                from pptx.util import Pt
                for r in title_shape.text_frame.paragraphs[1].runs:
                    base = r.font.size.pt if r.font.size else 40
                    r.font.size = Pt(max(15, int(base * 0.42)))
                    r.font.bold = False
            tc.shrink_text_to_shape(title_shape)
        else:
            # 封面无可用标题形状（标题是图片/无文字）→ 自绘标题块，绝不允许无标题封面
            rc.on_dark = bool((roles.get("cover") or {}).get("bg_is_dark"))
            if title_shape is not None:
                tc.set_text_keep_style(title_shape.text_frame, [""])
            _text(rc, slide, sd.title, page_w * 0.08, page_h * 0.58, page_w * 0.6, 1.25,
                  34, rc.heading_color, bold=True, font=rc.theme.font_title,
                  min_size=22, page=sd.page)
            if sd.subtitle:
                _text(rc, slide, sd.subtitle, page_w * 0.082, page_h * 0.58 + 1.3,
                      page_w * 0.55, 0.6, 15, rc.muted_color, page=sd.page)
        return

    if stype == "toc" and src_slide("toc") is not None:
        slide = tc.clone_slide(prs, src_slide("toc"))
        removed = tc.remove_ending_texts(slide)  # 目录页也不允许出现"谢谢"类残留
        if removed:
            rc.note(sd.page, "frame_cleaned", f"目录页清理残留文本 {removed} 处")
        items = _toc_items(data)
        mode, area = tc.prepare_toc(slide, items)
        if mode != "replaced":
            rc.on_dark = bool((roles.get("toc") or {}).get("bg_is_dark"))
            if area is None:
                area = content_area(roles.get("toc"))
            x, y, w, h = area
            # 区域夹紧到页面内并留出底部余量
            page_w = Emu(prs.slide_width).inches
            page_h = Emu(prs.slide_height).inches
            w = min(w, page_w - x - 0.5)
            h = min(max(h, 0.9 * len(items)), page_h - y - 0.5)
            BODY_PAINTERS["toc"](rc, slide, data, x, y, w, h)
        return

    if stype == "ending" and src_slide("ending") is not None:
        tc.clone_slide(prs, src_slide("ending"))  # 尾页原样保留
        return

    if stype == "section" and src_slide("section") is not None:
        slide = tc.clone_slide(prs, src_slide("section"))
        title_shape = tc.biggest_text_shape(slide)
        if title_shape is not None:
            tc.set_text_keep_style(title_shape.text_frame, [sd.title])
            tc.shrink_text_to_shape(title_shape)
        return

    # ---- 其余页面：拍平克隆内容框架页作底版（版式装饰保留、版式污染物过滤） ----
    if frame_info is not None:
        slide = tc.clone_slide_flat(prs, prs.slides[frame_info["slide_index"]], blank_layout)
        title_shape = tc.find_title_shape(slide)
        rc.on_dark = bool(frame_info.get("bg_is_dark"))  # 深底页自动切浅色文字
        x, y, w, h = content_area(frame_info)
        # 清理模板残留：示例文本/"谢谢"类文本/正文区大图（消除叠压与尾页内容泄漏）
        zone = (int(x * _EMU_IN), int(y * _EMU_IN),
                int((x + w) * _EMU_IN), int((y + h) * _EMU_IN))
        page_area = int(prs.slide_width or 0) * int(prs.slide_height or 0)
        cleaned = tc.clean_cloned_frame(slide, [title_shape], zone, page_area)
        # 标题带内的残留演示文字（旧标题/"汇报人："等）一并清除，避免与新标题叠压
        cleaned += tc.clean_title_band(slide, title_shape, int(y * _EMU_IN))
        if cleaned:
            rc.note(sd.page, "frame_cleaned", f"清理模板残留示例内容 {cleaned} 处")
        if stype == "section":
            if title_shape is not None:
                tc.set_text_keep_style(title_shape.text_frame, [""])
            BODY_PAINTERS["section"](rc, slide, data, x, y, w, h)
            return
        if title_shape is not None:
            tc.set_text_keep_style(title_shape.text_frame, [sd.title])
            # 细条标题带/继承字号场景下按形状适配字号，避免过小或溢出
            tc.shrink_text_to_shape(title_shape, min_size=13, default_max=24)
        else:  # 框架页没有标题形状 → 自绘标题（配色随背景自适应）
            _text(rc, slide, sd.title, x, 0.42, w, 0.7, 24, rc.heading_color,
                  bold=True, font=rc.theme.font_title, min_size=18, page=sd.page)
        if sd.subtitle:
            from app.ppt import text_engine
            # 副标题是辅助信息：单行放不下就整体省略，避免截断出现残句
            if text_engine.fits(sd.subtitle, w * 72, 0.42 * 72, 12.5):
                _text(rc, slide, sd.subtitle, x, y - 0.08, w, 0.4, 12.5, rc.muted_color,
                      page=sd.page)
                y += 0.42
                h -= 0.42
            else:
                rc.note(sd.page, "subtitle_skipped", "副标题过长已省略（避免截断残句）")
        body = BODY_PAINTERS.get(stype if stype != "ending" else "summary",
                                 BODY_PAINTERS["title_content"])
        body(rc, slide, data, x, y, w, h)
        return

    # 无内容框架页 → 模板空版式 + 系统整页绘制（仍保留母版背景）
    slide = prs.slides.add_slide(blank_layout)
    LAYOUT_PAINTERS.get(stype, paint_title_content)(rc, slide, data)


def _toc_items(data: dict) -> list[str]:
    for el in data.get("elements", []):
        if el.get("type") == "bullet_group":
            return [str(i) for i in el.get("items", [])]
    return []


def _emptiest_layout(prs):
    """选形状最少的版式做空白底（保留母版背景）。"""
    return min((lay for m in prs.slide_masters for lay in m.slide_layouts),
               key=lambda la: len(la.shapes))


# ================= 系统路径 =================

def _render_with_system(rc, spec, out_path: str) -> list[int]:
    degraded: list[int] = []
    try:
        prs = Presentation()
        prs.slide_width = Inches(PAGE_W)
        prs.slide_height = Inches(PAGE_H)
        blank = prs.slide_layouts[6]
    except Exception as e:
        logger.error("Presentation 初始化失败: %s", e)
        raise PPTError("E4001", str(e)) from e

    for sd in sorted(spec.slides, key=lambda s: s.page):
        slide = prs.slides.add_slide(blank)
        data = sd.model_dump()
        painter = LAYOUT_PAINTERS.get(sd.type, paint_title_content)
        try:
            painter(rc, slide, data)
        except Exception as e:
            logger.warning("第 %d 页 [%s] 绘制失败，降级为通用版式: %s", sd.page, sd.type, e)
            degraded.append(sd.page)
            rc.note(sd.page, "page_degraded", f"版式 {sd.type} 绘制失败已降级: {e}")
            try:
                fallback = {**data, "elements": [el for el in data.get("elements", [])
                                                 if el.get("type") == "bullet_group"]}
                if not fallback["elements"]:
                    fallback["elements"] = [{"type": "bullet_group",
                                             "items": [data.get("key_message") or data.get("title", "")]}]
                paint_title_content(rc, slide, fallback)
            except Exception as e2:
                logger.error("第 %d 页降级绘制仍失败(保留空页): %s", sd.page, e2)
    _save(prs, out_path)
    return degraded


def _save(prs, out_path: str) -> None:
    try:
        prs.save(out_path)
    except Exception as e:
        logger.error("PPTX 保存失败: %s", e)
        raise PPTError("E4001", f"保存失败: {e}") from e
