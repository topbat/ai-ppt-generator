"""AI 模板生成器：按八个业务维度参数化生成 PPT 模板（确定性、零 LLM 成本）。

八维参数：行业 / 用户群 / 风格 / 数据内容 / 主题 / 关键字 / 国家 / 季节。
映射规则：行业→基础配色；风格→构图与明暗；季节→强调色调温；国家→装饰语言；
用户群→字号与留白档位；数据内容→内容页组合（5~8 页）；主题/关键字→封面占位文案与命名。

产出结构（硬性要求）：封面 + 目录页 + 章节页 + 5~8 个内容页 + 尾页；
命名一律以 "AI" 为前缀。页面设计与模板解析器的角色体系对齐，生成任务可直接克隆复用。

设计约束（勿破坏）：封面/目录/章节页的可替换文本结构是克隆链路的锚点——
封面只有最大文本框会被替换、目录按顶层文本条目原位替换（多余条目会被删除）；
因此视觉装饰一律用无文字形状（矩形/圆环），不得给结构页添加多余装饰文字。
"""
import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.logging import get_logger

logger = get_logger(__name__)

W, H = 13.333, 7.5
_FONT = "Microsoft YaHei"

# ---- 八维选项与映射表（对外经 /templates/ai-options 提供） ----

INDUSTRIES = {  # 行业 → (primary, secondary, accent)
    "通用": ("1B3A6B", "0F2547", "2F80ED"),
    "科技互联网": ("4C3FD6", "241E6B", "7C6CFF"),
    "金融投资": ("103C64", "081F36", "C9A227"),
    "政务机关": ("B02A24", "6E1B17", "D4AC0D"),
    "医疗健康": ("0E7C6B", "07443B", "2EC4B6"),
    "教育培训": ("D97706", "8A4B04", "2F80ED"),
    "制造能源": ("35506B", "1F3143", "E67E22"),
    "地产建筑": ("1E5A46", "0F2E23", "B08D57"),
    "文旅消费": ("C2452D", "742A1B", "F0A202"),
    "咨询服务": ("2F4550", "1B2A32", "586F7C"),
}

AUDIENCES = ["高管汇报", "客户提案", "内部团队", "公众演讲", "学术评审"]
STYLES = ["商务稳重", "简约留白", "科技暗黑", "活力明快", "庄重大气"]
DATA_CONTENTS = ["均衡混合", "数据图表型", "文字要点型", "图文展示型"]
COUNTRIES = ["中国", "国际通用", "日韩", "欧美"]
SEASONS = ["不限", "春", "夏", "秋", "冬"]

# 风格 → (封面构图, 是否深底, 底色覆盖)
_STYLE_MAP = {
    "商务稳重": ("side", False, "FFFFFF"),
    "简约留白": ("center", False, "FAFAFA"),
    "科技暗黑": ("center", True, "12141E"),
    "活力明快": ("band", False, "FFFFFF"),
    "庄重大气": ("band", False, "FFFFFF"),
}
# 季节 → 强调色调温（向目标色相混合比例）
_SEASON_TINT = {"春": ("3CB371", 0.35), "夏": ("1E90FF", 0.3),
                "秋": ("D2691E", 0.35), "冬": ("4682B4", 0.3)}
# 用户群 → (封面标题字号, 正文留白档)
_AUDIENCE_SCALE = {"高管汇报": (44, 1.15), "客户提案": (40, 1.0), "内部团队": (36, 0.9),
                   "公众演讲": (46, 1.2), "学术评审": (38, 1.0)}
# 数据内容 → 内容页组合（5~8 页，取自九种内容页型；表格/图表为原生对象，解析器可识别）
_CONTENT_SETS = {
    "均衡混合": ["frame", "cards", "chart", "timeline", "two_col", "table", "numbers"],
    "数据图表型": ["frame", "chart", "numbers", "table", "timeline", "two_col"],
    "文字要点型": ["frame", "two_col", "cards", "timeline", "arch", "table"],
    "图文展示型": ["frame", "cards", "image", "timeline", "numbers", "arch", "chart"],
}

# 启动播种的 10 套预设（空库时自动生成）
SEED_PRESETS = [
    {"industry": "通用", "style": "商务稳重", "data_content": "均衡混合", "theme": "商务汇报"},
    {"industry": "科技互联网", "style": "科技暗黑", "data_content": "数据图表型", "theme": "科技发布"},
    {"industry": "金融投资", "style": "庄重大气", "data_content": "数据图表型", "theme": "金融分析"},
    {"industry": "政务机关", "style": "庄重大气", "data_content": "文字要点型", "theme": "政务汇报"},
    {"industry": "医疗健康", "style": "简约留白", "data_content": "均衡混合", "theme": "医疗方案"},
    {"industry": "教育培训", "style": "活力明快", "data_content": "图文展示型", "theme": "教育课件"},
    {"industry": "制造能源", "style": "商务稳重", "data_content": "数据图表型", "theme": "产业报告"},
    {"industry": "地产建筑", "style": "商务稳重", "data_content": "均衡混合", "theme": "项目介绍"},
    {"industry": "文旅消费", "style": "活力明快", "data_content": "图文展示型",
     "theme": "文旅推介", "country": "中国", "season": "秋"},
    {"industry": "咨询服务", "style": "简约留白", "data_content": "文字要点型", "theme": "咨询方案"},
]


def ai_options() -> dict:
    """八维参数可选项（主题/关键字为自由文本，不在此列）。"""
    return {"industries": list(INDUSTRIES), "audiences": AUDIENCES, "styles": STYLES,
            "data_contents": DATA_CONTENTS, "countries": COUNTRIES, "seasons": SEASONS}


# ---- 颜色工具 ----

def _mix(hex_a: str, hex_b: str, ratio: float) -> str:
    """a 向 b 混合 ratio（0~1）。"""
    ra, ga, ba = (int(hex_a[i:i + 2], 16) for i in (0, 2, 4))
    rb, gb, bb = (int(hex_b[i:i + 2], 16) for i in (0, 2, 4))
    return "".join(f"{int(x + (y - x) * ratio):02X}" for x, y in ((ra, rb), (ga, gb), (ba, bb)))


def _tint(hex_color: str, factor: float) -> str:
    return _mix(hex_color, "FFFFFF", factor)


def normalize_params(params: dict | None) -> dict:
    """参数规整：非法值回退默认，保证生成器永不因输入炸掉。"""
    p = dict(params or {})
    p["industry"] = p.get("industry") if p.get("industry") in INDUSTRIES else "通用"
    p["audience"] = p.get("audience") if p.get("audience") in _AUDIENCE_SCALE else "客户提案"
    p["style"] = p.get("style") if p.get("style") in _STYLE_MAP else "商务稳重"
    p["data_content"] = (p.get("data_content")
                         if p.get("data_content") in _CONTENT_SETS else "均衡混合")
    p["country"] = p.get("country") if p.get("country") in COUNTRIES else "国际通用"
    p["season"] = p.get("season") if p.get("season") in SEASONS else "不限"
    p["theme"] = str(p.get("theme") or "").strip()[:20]
    p["keywords"] = str(p.get("keywords") or "").strip()[:40]
    return p


def template_name(params: dict) -> str:
    """AI 模板命名：一律以 AI 为前缀。"""
    p = normalize_params(params)
    core = p["theme"] or p["industry"]
    return f"AI·{core}·{p['style']}"[:60]


def build_ai_template(params: dict | None = None) -> bytes:
    """按八维参数生成模板，返回 pptx 字节。

    结构：封面 → 目录 → 章节页 → 内容页 5~8 → 尾页。"""
    p = normalize_params(params)
    primary, secondary, accent = INDUSTRIES[p["industry"]]
    cover_kind, dark, bg = _STYLE_MAP[p["style"]]

    # 季节调温：强调色向季节色相混合
    if p["season"] in _SEASON_TINT:
        target, ratio = _SEASON_TINT[p["season"]]
        accent = _mix(accent, target, ratio)
    # 国家装饰语言：中国→装饰线偏暖金；日韩→低饱和；欧美→偏冷灰
    deco = accent
    if p["country"] == "中国":
        deco = _mix(accent, "D4AC0D", 0.5)
    elif p["country"] == "日韩":
        deco = _mix(accent, "FFFFFF", 0.35)
    elif p["country"] == "欧美":
        deco = _mix(accent, "6B7A8F", 0.3)

    title_size, _space = _AUDIENCE_SCALE[p["audience"]]
    text_main = "FFFFFF" if dark else secondary
    text_sub = "C8D0DC" if dark else "8A94A6"
    band_color = secondary if dark else primary
    content_pages = _CONTENT_SETS[p["data_content"]]

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def slide():
        sl = prs.slides.add_slide(blank)
        _rect(sl, 0, 0, W, H, bg)
        return sl

    def _rect(sl, x, y, w, h, color, rounded=False, line=None):
        shp = sl.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(color)
        if line:
            shp.line.color.rgb = RGBColor.from_string(line)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _oval(sl, x, y, w, h, fill=None, line=None, line_w=1.1):
        """圆/圆环装饰：fill=None 即空心环（仅描边），全部无文字，不干扰克隆链路。"""
        shp = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor.from_string(fill)
        else:
            shp.fill.background()
        if line:
            shp.line.color.rgb = RGBColor.from_string(line)
            shp.line.width = Pt(line_w)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _text(sl, txt, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = _FONT
        run.font.color.rgb = RGBColor.from_string(color)
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": _FONT}))
        return tb

    def _cell(cell, txt, size, color, bold=False, fill=None, align=PP_ALIGN.LEFT):
        """表格单元格：填充 + 中英文字体统一设置。"""
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(fill)
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = _FONT
        run.font.color.rgb = RGBColor.from_string(color)
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": _FONT}))

    def title_band(sl, label="页面标题"):
        """内容页统一标题带（风格随参数）。

        页码占位用 "01/10"：纯数字 "01" 会命中解析器的章节页判定（大数字+短文本），
        导致内容框架页被误分类。"""
        _rect(sl, 0, 0, W, 1.15, band_color)
        _rect(sl, 0, 1.15, W, 0.06, deco)
        _rect(sl, 0.55, 0.38, 0.13, 0.42, deco)   # 标题前置强调块
        _text(sl, label, 0.9, 0.28, 10.5, 0.65, 24, "FFFFFF", bold=True)
        _rect(sl, 0.9, H - 0.42, 1.6, 0.03, _tint(primary, 0.8) if not dark
              else _mix(secondary, "FFFFFF", 0.2))  # 页脚装饰线
        _text(sl, "01/10", W - 1.3, H - 0.5, 0.9, 0.35, 10, text_sub, align=PP_ALIGN.RIGHT)

    cover_title = p["theme"] or "演示文稿标题"
    cover_sub = p["keywords"] or "副标题 · 汇报人"

    # ---- 1 封面（装饰一律用无文字形状，仅最大文本框会被生成链路替换）----
    sl = slide()
    ring = _mix(primary, "FFFFFF", 0.30)   # 色块上的空心环描边色
    if cover_kind == "side":
        _rect(sl, W * 0.62, 0, W * 0.38, H, primary)
        _rect(sl, W * 0.60, 0, 0.10, H, deco)
        # 右侧色块：同心圆环 + 强调圆点，克制但有记忆点
        _oval(sl, W * 0.70, 0.95, 2.7, 2.7, line=ring, line_w=1.3)
        _oval(sl, W * 0.74, 1.5, 1.6, 1.6, line=_mix(primary, "FFFFFF", 0.5), line_w=1.0)
        _oval(sl, W * 0.79, 4.9, 0.5, 0.5, fill=deco)
        _oval(sl, W * 0.70, 5.9, 0.24, 0.24, fill=_mix(primary, "FFFFFF", 0.45))
        _text(sl, cover_title, 0.9, 2.7, 6.6, 1.4, title_size, text_main, bold=True)
        _rect(sl, 0.95, 2.45, 1.5, 0.06, accent)
        _text(sl, cover_sub, 0.95, 4.2, 6.0, 0.6, 17, text_sub)
        _rect(sl, 0.95, H - 0.8, 2.1, 0.04, _tint(primary, 0.75))  # 页脚定位线
    elif cover_kind == "band":
        _oval(sl, W - 3.1, 0.5, 2.0, 2.0, line=_tint(primary, 0.72), line_w=1.2)
        _oval(sl, 0.7, H - 2.3, 1.3, 1.3, line=_tint(primary, 0.78), line_w=1.0)
        _rect(sl, 0, H * 0.30 - 0.08, 2.6, 0.08, deco)   # 主带上沿角标线
        _rect(sl, 0, H * 0.30, W, H * 0.34, primary)
        _rect(sl, 0, H * 0.64, W, 0.08, deco)
        _oval(sl, W - 2.6, H * 0.36, 1.5, 1.5, line=ring, line_w=1.1)
        _text(sl, cover_title, 1.2, H * 0.36, W - 2.4, 1.3, title_size, "FFFFFF",
              bold=True, align=PP_ALIGN.CENTER)
        _text(sl, cover_sub, 1.2, H * 0.70, W - 2.4, 0.6, 16, text_sub, align=PP_ALIGN.CENTER)
    else:  # center
        _rect(sl, 0, 0, W, H, secondary if dark else bg)
        edge = _mix(secondary, "FFFFFF", 0.16) if dark else _tint(primary, 0.82)
        _oval(sl, W - 3.4, -1.2, 3.6, 3.6, line=edge, line_w=1.3)   # 右上出血圆环
        _oval(sl, -1.0, H - 2.4, 2.6, 2.6, line=edge, line_w=1.1)   # 左下出血圆环
        _oval(sl, W - 1.7, 2.1, 0.3, 0.3, fill=accent)
        _rect(sl, W / 2 - 1.0, 2.15, 2.0, 0.07, accent)
        _text(sl, cover_title, 1.2, 2.6, W - 2.4, 1.4, title_size + 2,
              "FFFFFF" if dark else primary, bold=True, align=PP_ALIGN.CENTER)
        _text(sl, cover_sub, 1.2, 4.3, W - 2.4, 0.6, 16, text_sub, align=PP_ALIGN.CENTER)
        _rect(sl, W / 2 - 0.35, 5.15, 0.7, 0.045, deco)
        _rect(sl, 0, H - 0.5, W, 0.5, primary)

    # ---- 2 目录（双列排版；条目保持"单个顶层文本框"结构，克隆时原位替换/裁剪）----
    sl = slide()
    _rect(sl, 0, 0, 0.5, H, primary)
    _rect(sl, 0.5, 0, 0.06, H, deco)
    _oval(sl, W - 2.3, -0.9, 2.6, 2.6, line=_tint(primary, 0.8) if not dark
          else _mix(secondary, "FFFFFF", 0.16), line_w=1.2)
    _text(sl, "目录", 1.1, 0.72, 3.0, 1.0, 34, text_main, bold=True)
    _text(sl, "CONTENTS", 1.13, 1.62, 3.5, 0.4, 11, text_sub)  # 克隆时按目录标题词排除
    _rect(sl, 1.15, 2.1, 1.1, 0.06, deco)
    for i in range(6):
        col, row = i % 2, i // 2
        _text(sl, f"0{i + 1}   章节标题占位", 1.35 + col * 5.9, 2.75 + row * 1.25,
              5.4, 0.6, 17, text_main, bold=True)

    # ---- 3 章节过渡页 ----
    sl = slide()
    _rect(sl, 0, 0, W, H, primary)
    _oval(sl, W - 3.6, H - 3.4, 3.0, 3.0, line=ring, line_w=1.4)
    _oval(sl, W - 4.1, H - 1.3, 0.4, 0.4, fill=deco)
    _rect(sl, W - 1.0, 0.65, 0.35, 0.06, deco)
    _text(sl, "01", 1.1, 2.0, 3.4, 1.9, 88, _tint(primary, 0.55), bold=True)
    _rect(sl, 1.2, 4.0, 1.6, 0.06, "FFFFFF")
    _text(sl, "章节标题", 1.15, 4.25, 9.5, 1.0, 32, "FFFFFF", bold=True)

    # ---- 4~N 内容页（5~8 页，按数据内容组合） ----
    card_bg = _mix(bg, "FFFFFF", 0.5) if dark else _tint(primary, 0.94)
    body_text = "E8ECF4" if dark else secondary
    shadow = _mix(bg, "000000", 0.35) if dark else _mix("FFFFFF", secondary, 0.12)
    hair = _mix(secondary, "FFFFFF", 0.2) if dark else _tint(primary, 0.75)  # 细描边色

    for kind in content_pages:
        sl = slide()
        if kind == "frame":       # 内容框架页：标题带 + 干净正文区（克隆底版首选）
            title_band(sl)
        elif kind == "cards":     # 三卡片页（底层偏移色块模拟柔和投影）
            title_band(sl)
            cw = (W - 1.8 - 0.8) / 3
            for i in range(3):
                cx = 0.9 + i * (cw + 0.4)
                _rect(sl, cx + 0.06, 2.08, cw, 3.6, shadow, rounded=True)
                _rect(sl, cx, 2.0, cw, 3.6, card_bg, rounded=True, line=hair)
                _rect(sl, cx + 0.02, 2.0, cw - 0.04, 0.08, accent)
                _text(sl, f"要点标题{i + 1}", cx + 0.3, 2.35, cw - 0.6, 0.5, 16,
                      body_text, bold=True)
                _text(sl, "在此填写要点说明文字，概述核心观点与支撑信息。",
                      cx + 0.3, 3.0, cw - 0.6, 2.2, 12, text_sub)
        elif kind == "chart":     # 图表页：原生柱状图（可编辑）
            title_band(sl, "数据图表标题")
            data = CategoryChartData()
            data.categories = ["2024", "2025", "2026"]
            data.add_series("示例数据", (82, 116, 158))
            gf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                     Inches(2.2), Inches(1.7), Inches(9.0), Inches(4.9), data)
            chart = gf.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            try:
                chart.plots[0].series[0].format.fill.solid()
                chart.plots[0].series[0].format.fill.fore_color.rgb = \
                    RGBColor.from_string(accent)
            except Exception:
                pass
        elif kind == "two_col":   # 两栏页
            title_band(sl, "两栏对比标题")
            for i in range(2):
                half = (W - 1.8 - 0.5) / 2
                cx = 0.9 + i * (half + 0.5)
                _rect(sl, cx, 1.75, half, 0.5, _tint(primary, 0.85) if not dark
                      else _mix(secondary, "FFFFFF", 0.12))
                _text(sl, f"栏目标题{i + 1}", cx + 0.2, 1.8, half - 0.4, 0.4, 15,
                      body_text, bold=True)
                _text(sl, "▪ 要点一说明文字占位\n▪ 要点二说明文字占位\n▪ 要点三说明文字占位",
                      cx + 0.2, 2.5, half - 0.4, 3.4, 13, text_sub)
        elif kind == "numbers":   # 大数字页
            title_band(sl, "关键指标标题")
            vals = [("85%", "指标一说明文字占位"), ("1.2亿", "指标二说明文字占位"),
                    ("3年", "指标三说明文字占位")]
            cw = (W - 1.8) / 3
            for i, (v, label) in enumerate(vals):
                cx = 0.9 + i * cw
                _rect(sl, cx + cw / 2 - 0.3, 2.42, 0.6, 0.05, accent)  # 数值上方刻度线
                _text(sl, v, cx, 2.6, cw, 1.3, 54, accent, bold=True, align=PP_ALIGN.CENTER)
                _text(sl, label, cx, 4.1, cw, 0.6, 14, text_sub, align=PP_ALIGN.CENTER)
                if i:
                    _rect(sl, cx, 2.9, 0.014, 1.6, hair)
        elif kind == "image":     # 图文页（左图占位 + 右文）
            title_band(sl, "图文页标题")
            _rect(sl, 0.96, 1.98, 5.6, 4.4, shadow, rounded=True)
            _rect(sl, 0.9, 1.9, 5.6, 4.4, card_bg, rounded=True, line=hair)
            _text(sl, "图片占位", 0.9, 3.8, 5.6, 0.6, 14, text_sub, align=PP_ALIGN.CENTER)
            _rect(sl, 6.9, 2.35, 0.45, 0.05, accent)
            _text(sl, "▪ 图注要点一占位\n▪ 图注要点二占位\n▪ 图注要点三占位",
                  6.9, 2.6, 5.4, 3.2, 14, body_text)
        elif kind == "timeline":  # 时间轴/推进计划页：主轴 + 节点圆 + 上下错落文案
            title_band(sl, "推进计划标题")
            steps = 4
            seg = (W - 1.8) / steps
            axis_y = 4.0
            _rect(sl, 0.9, axis_y + 0.10, W - 1.8, 0.035, hair)
            for i in range(steps):
                cx = 0.9 + seg * i + seg / 2
                _oval(sl, cx - 0.13, axis_y, 0.26, 0.26,
                      fill=accent if i == 0 else bg, line=accent, line_w=1.4)
                _text(sl, f"阶段{i + 1}", cx - seg / 2 + 0.2, axis_y - 0.85,
                      seg - 0.4, 0.45, 15, body_text, bold=True, align=PP_ALIGN.CENTER)
                _text(sl, "阶段说明文字占位", cx - seg / 2 + 0.25, axis_y + 0.55,
                      seg - 0.5, 1.0, 12, text_sub, align=PP_ALIGN.CENTER)
        elif kind == "table":     # 表格页：原生表格（解析器识别为 table 版式）
            title_band(sl, "数据表格标题")
            rows, cols = 4, 4
            tbl = sl.shapes.add_table(rows, cols, Inches(1.0), Inches(1.95),
                                      Inches(W - 2.0), Inches(3.5)).table
            headers = ["对比维度", "指标一", "指标二", "指标三"]
            row_a = "FFFFFF" if not dark else _mix(secondary, "FFFFFF", 0.08)
            row_b = _tint(primary, 0.95) if not dark else _mix(secondary, "FFFFFF", 0.14)
            for c, htxt in enumerate(headers):
                _cell(tbl.cell(0, c), htxt, 13, "FFFFFF", bold=True, fill=band_color,
                      align=PP_ALIGN.CENTER)
            for r in range(1, rows):
                fill = row_a if r % 2 else row_b
                _cell(tbl.cell(r, 0), f"项目{r}占位", 12, body_text, bold=True, fill=fill)
                for c in range(1, cols):
                    _cell(tbl.cell(r, c), "内容占位", 12, body_text, fill=fill,
                          align=PP_ALIGN.CENTER)
        elif kind == "arch":      # 架构分层页：应用/平台/基础设施三层
            title_band(sl, "总体架构标题")
            lx, lw = 1.5, W - 3.0
            _rect(sl, lx, 1.85, lw, 0.7, band_color, rounded=True)
            _text(sl, "应用层占位", lx, 2.0, lw, 0.45, 15, "FFFFFF", bold=True,
                  align=PP_ALIGN.CENTER)
            mw = (lw - 0.5) / 3
            for i in range(3):
                mx = lx + i * (mw + 0.25)
                _rect(sl, mx, 2.8, mw, 1.4, card_bg, rounded=True, line=hair)
                _rect(sl, mx + 0.02, 2.8, mw - 0.04, 0.07, accent)
                _text(sl, f"能力模块{i + 1}占位", mx, 3.3, mw, 0.5, 13, body_text,
                      bold=True, align=PP_ALIGN.CENTER)
            _rect(sl, lx, 4.45, lw, 0.7, _mix(primary, secondary, 0.4), rounded=True)
            _text(sl, "数据与平台层占位", lx, 4.6, lw, 0.45, 14, "FFFFFF",
                  align=PP_ALIGN.CENTER)
            _rect(sl, lx, 5.4, lw, 0.7, secondary, rounded=True)
            _text(sl, "基础设施层占位", lx, 5.55, lw, 0.45, 14, "FFFFFF",
                  align=PP_ALIGN.CENTER)

    # ---- 尾页 ----
    sl = slide()
    end_bg = secondary if dark else primary
    _rect(sl, 0, 0, W, H, end_bg)
    end_ring = _mix(end_bg, "FFFFFF", 0.25)
    _oval(sl, -1.1, -1.1, 3.2, 3.2, line=end_ring, line_w=1.3)   # 左上出血圆环
    _oval(sl, W - 2.4, H - 2.4, 3.2, 3.2, line=end_ring, line_w=1.3)   # 右下出血圆环
    _oval(sl, W - 3.0, H - 1.1, 0.32, 0.32, fill=deco)
    _text(sl, "感谢聆听", 1.2, 2.9, W - 2.4, 1.4, 46, "FFFFFF", bold=True,
          align=PP_ALIGN.CENTER)
    _rect(sl, W / 2 - 0.9, 4.45, 1.8, 0.06, deco)

    buf = io.BytesIO()
    prs.save(buf)
    logger.info("AI 模板生成完成：%s（行业=%s 风格=%s 数据=%s 国家=%s 季节=%s，内容页 %d 个，共 %d 页）",
                template_name(p), p["industry"], p["style"], p["data_content"],
                p["country"], p["season"], len(content_pages), 4 + len(content_pages))
    return buf.getvalue()


# ---- 兼容保留：旧风格清单接口（新代码请使用 ai_options / build_ai_template） ----
GALLERY_STYLES = [{"name": s, "primary": INDUSTRIES["通用"][0], "dark": _STYLE_MAP[s][1],
                   "cover": _STYLE_MAP[s][0]} for s in STYLES]


def build_template(style: dict) -> bytes:
    """旧接口兼容：按风格名生成（映射到新参数化生成器）。"""
    return build_ai_template({"style": style.get("name", "商务稳重")})
