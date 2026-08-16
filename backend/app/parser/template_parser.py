"""PPT 模板解析器：提取 Design Token + 逐页角色分类（供克隆引擎复用模板页）。

角色体系（template_slides.slide_type）：
- cover         封面页（克隆后替换标题/副标题）
- toc           目录页（克隆后替换条目文字）
- content_frame 内容框架页（仅标题+装饰，正文区干净——正文页的克隆底版）
- ending        尾页（"感谢聆听"类，原样克隆）
- section       章节过渡页
- 其余沿用内容类型枚举（three_cards/bar_chart/...，供参考评分）
layout_meta 记录标题形状几何（克隆后正文绘制区依此计算）。
PARSER_VERSION 变更时 PARSE_TPL 阶段自动触发重新解析。
"""
import re
import zipfile

from pptx import Presentation
from pptx.util import Emu

from app.core.errors import PPTError
from app.core.logging import get_logger

logger = get_logger(__name__)

PARSER_VERSION = 5  # 分类规则变更时 +1，旧解析结果自动失效重解析（v5: 版式层尾页识别）

_TOC_WORDS = ("目录", "目 录", "contents", "content", "agenda", "provide")
_ENDING_WORDS = ("感谢", "谢谢", "thank", "thanks", "the end", "再见")


def parse_template(path: str) -> dict:
    try:
        prs = Presentation(path)
    except Exception as e:
        logger.error("模板打开失败: %s", e)
        raise PPTError("E2004", str(e)) from e

    tokens = _extract_theme_tokens(path)
    tokens["page_width_emu"] = int(prs.slide_width or Emu(12192000))
    tokens["page_height_emu"] = int(prs.slide_height or Emu(6858000))

    # 字体采收：模板页面实际使用的中文字体优先于主题声明（字体选择与模板一致）
    harvested = _harvest_fonts(prs)
    if harvested:
        tokens["font_title"] = harvested
        tokens["font_body"] = harvested
        logger.info("采收模板实际用字：%s", harvested)

    # 主色采收：模板页面实际使用的高频饱和色优先于主题 accent1
    # （大量模板主题色仍是 Office 默认蓝，但页面实际是红/金等品牌色，
    #  不采收会导致章节序号/强调色与模板配色冲突）
    harvested_color = _harvest_primary_color(prs)
    if harvested_color:
        tokens["primary"] = harvested_color
        logger.info("采收模板实际主色：#%s", harvested_color)

    slides = []
    scheme = tokens.get("scheme") or {}
    for idx, slide in enumerate(prs.slides):
        stype, conf, meta = _classify_slide(slide, idx, len(prs.slides), prs)
        meta["parser_ver"] = PARSER_VERSION
        # 背景色估计：深底页正文自动切浅色文字（字色与背景协调的关键）
        bg = _estimate_bg_color(slide, prs, scheme)
        meta["bg_color"] = bg
        meta["bg_is_dark"] = _luma(bg) < 140
        slides.append({"slide_index": idx, "slide_type": stype, "confidence": conf,
                       "layout_meta": meta})

    if not slides:
        logger.warning("模板不含任何页面，将仅使用其主题配色与字体")
    roles = [s["slide_type"] for s in slides]
    logger.info("模板解析完成(v%d)：%d 页 角色=%s 主色=#%s 标题字体=%s",
                PARSER_VERSION, len(slides), roles, tokens.get("primary"), tokens.get("font_title"))
    return {"design_tokens": tokens, "slide_count": len(slides), "slides": slides}


def _extract_theme_tokens(path: str) -> dict:
    """直接读 theme1.xml 提取主题色与字体（python-pptx 未暴露完整主题 API）。"""
    tokens = {"primary": "1677FF", "secondary": "13234E", "accent": "36CFC9",
              "text": "1F1F1F", "background": "FFFFFF",
              "font_title": "Microsoft YaHei", "font_body": "Microsoft YaHei"}
    try:
        with zipfile.ZipFile(path) as z:
            names = [n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml", n)]
            if not names:
                return tokens
            xml = z.read(sorted(names)[0]).decode("utf-8", errors="ignore")
        colors = dict(re.findall(
            r'<a:(dk1|lt1|dk2|lt2|accent1|accent2|accent3)>.*?val="([0-9A-Fa-f]{6})"', xml, re.S))
        if colors.get("accent1"):
            tokens["primary"] = colors["accent1"].upper()
        if colors.get("accent2"):
            tokens["accent"] = colors["accent2"].upper()
        if colors.get("dk2"):
            tokens["secondary"] = colors["dk2"].upper()
        major = re.search(r"<a:majorFont>(.*?)</a:majorFont>", xml, re.S)
        minor = re.search(r"<a:minorFont>(.*?)</a:minorFont>", xml, re.S)
        for block, key in ((major, "font_title"), (minor, "font_body")):
            if not block:
                continue
            ea = re.search(r'<a:ea typeface="([^"]+)"', block.group(1))
            latin = re.search(r'<a:latin typeface="([^"]+)"', block.group(1))
            face = (ea and ea.group(1)) or (latin and latin.group(1))
            if face and not face.startswith("+"):
                tokens[key] = face
        # 完整配色方案表（背景色 schemeClr 解析用；含 sysClr 的 lastClr）
        scheme = dict(re.findall(
            r'<a:(dk1|lt1|dk2|lt2|accent[1-6]|hlink|folHlink)>\s*'
            r'<a:(?:srgbClr val|sysClr[^>]*?lastClr)="([0-9A-Fa-f]{6})"', xml))
        tokens["scheme"] = {k: v.upper() for k, v in scheme.items()}
    except Exception as e:
        logger.warning("主题 Token 提取失败，使用默认 Token(兜底): %s", e)
    return tokens


def _harvest_primary_color(prs) -> str | None:
    """统计模板页面实际使用的高频饱和填充色（排除灰阶/近白/近黑），作为主色。"""
    from collections import Counter
    counter: Counter = Counter()
    try:
        for slide in prs.slides:
            xml = slide._element.xml
            for c in re.findall(r'<a:solidFill>\s*<a:srgbClr val="([0-9A-Fa-f]{6})"', xml):
                c = c.upper()
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                if max(r, g, b) - min(r, g, b) < 45:
                    continue  # 灰阶/低饱和不作主色
                luma = 0.299 * r + 0.587 * g + 0.114 * b
                if luma < 30 or luma > 225:
                    continue
                counter[c] += 1
        best = counter.most_common(1)
        return best[0][0] if best and best[0][1] >= 3 else None
    except Exception:
        return None


def _harvest_fonts(prs) -> str | None:
    """统计模板页面 run 级实际使用的东亚字体，返回出现频次最高者。

    很多模板主题声明是 Calibri 等西文字体，但页面实际用微软雅黑/思源黑体，
    按主题取字会导致生成内容与模板字体不一致。"""
    from collections import Counter
    counter: Counter = Counter()
    try:
        for slide in prs.slides:
            xml = slide._element.xml
            for face in re.findall(r'<a:ea typeface="([^"]+)"', xml):
                if not face.startswith("+"):
                    counter[face] += 1
        best = counter.most_common(1)
        return best[0][0] if best else None
    except Exception:
        return None


# ================= 背景色估计 =================

def _luma(hex_color: str) -> float:
    try:
        r, g, b = (int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
        return 0.299 * r + 0.587 * g + 0.114 * b
    except Exception:
        return 255.0


def _xml_fill_color(xml: str, scheme: dict) -> str | None:
    """从 XML 片段解析首个填充颜色（srgbClr 或 schemeClr 映射到主题色）。"""
    m = re.search(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', xml)
    if m:
        return m.group(1).upper()
    m = re.search(r'<a:schemeClr val="(\w+)"', xml)
    if m:
        name = m.group(1)
        # bg1/tx1 等别名 → 主题方案色
        alias = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}.get(name, name)
        return scheme.get(alias) or {"lt1": "FFFFFF", "lt2": "F2F2F2",
                                     "dk1": "000000", "dk2": "1F1F1F"}.get(alias)
    return None


def _estimate_bg_color(slide, prs, scheme: dict) -> str:
    """估计页面主背景色：页背景 → 版式背景 → 母版背景 → 大面积覆盖形状。

    背景为图片时无法取色，按浅色处理（保守策略：深底误判为浅底只是不够醒目，
    浅底误判为深底会白字不可见，故宁判浅）。
    """
    from lxml import etree

    def bg_of(element) -> str | None:
        bg = element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
        if bg is None:
            return None
        bg = bg.find("{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
        if bg is None:
            return None
        return _xml_fill_color(etree.tostring(bg, encoding="unicode"), scheme)

    page_area = (prs.slide_width or 1) * (prs.slide_height or 1)

    # 1) 页面级显式背景
    color = bg_of(slide._element)
    if color:
        return color

    # 2) 整页图片背景：采样图片真实平均色（图片底模板的字色协调关键）
    for shp in slide.shapes:
        try:
            if (shp.shape_type == 13 and shp.width and shp.height
                    and shp.width * shp.height > page_area * 0.6):
                from io import BytesIO

                from PIL import Image
                img = Image.open(BytesIO(shp.image.blob)).convert("RGB")
                img.thumbnail((24, 24))
                px = list(img.getdata())
                n = max(1, len(px))
                r = sum(p[0] for p in px) // n
                g = sum(p[1] for p in px) // n
                b = sum(p[2] for p in px) // n
                return f"{r:02X}{g:02X}{b:02X}"
        except Exception:
            continue

    # 3) 大面积覆盖的纯色形状（如整页色块打底）
    for shp in slide.shapes:
        try:
            if (shp.shape_type in (1, 17) and shp.width and shp.height
                    and shp.width * shp.height > page_area * 0.55):
                from lxml import etree as _et
                color = _xml_fill_color(_et.tostring(shp._element, encoding="unicode"), scheme)
                if color:
                    return color
        except Exception:
            continue

    # 4) 版式/母版背景
    for element in (slide.slide_layout._element, slide.slide_layout.slide_master._element):
        color = bg_of(element)
        if color:
            return color
    return "FFFFFF"


def _slide_texts(slide) -> list[str]:
    """递归收集页面全部文本（穿透 Group）。"""
    texts = []

    def walk(shapes):
        for s in shapes:
            if s.shape_type == 6:  # GROUP
                walk(s.shapes)
            elif getattr(s, "has_text_frame", False):
                t = s.text_frame.text.strip()
                if t:
                    texts.append(t)
    walk(slide.shapes)
    return texts


def _title_geometry(slide) -> dict | None:
    """记录标题形状几何（EMU），克隆后据此计算正文绘制区。"""
    from app.ppt.template_clone import find_title_shape
    shp = find_title_shape(slide)
    if shp is None or shp.top is None:
        return None
    return {"left": int(shp.left or 0), "top": int(shp.top or 0),
            "width": int(shp.width or 0), "height": int(shp.height or 0)}


def _classify_slide(slide, idx: int, total: int, prs) -> tuple[str, float, dict]:
    shapes = list(slide.shapes)
    texts = _slide_texts(slide)
    all_text = " ".join(texts)
    pictures = [s for s in shapes if s.shape_type == 13]
    has_chart = any(getattr(s, "has_chart", False) for s in shapes)
    has_table = any(getattr(s, "has_table", False) for s in shapes)
    # 有效正文字符数（排除页码类）
    body_chars = sum(len(t) for t in texts if not re.fullmatch(r"[\d\s/]*", t))
    meta = {"shape_count": len(shapes), "text_count": len(texts),
            "picture_count": len(pictures), "has_chart": has_chart, "has_table": has_table,
            "title_geo": _title_geometry(slide)}

    # 强文本信号优先（尾页判定包含版式层文字：谢谢艺术字常放在版式里）
    layout_text = ""
    try:
        def walk_layout(shapes):
            out = []
            for s in shapes:
                if s.shape_type == 6:
                    out.extend(walk_layout(s.shapes))
                elif getattr(s, "has_text_frame", False):
                    out.append(s.text_frame.text)
            return out
        layout_text = " ".join(walk_layout(slide.slide_layout.shapes))
    except Exception:
        pass
    if any(w in (all_text + " " + layout_text).lower() for w in _ENDING_WORDS) and body_chars < 40:
        return "ending", 0.95, meta
    if any(w in all_text.lower() for w in _TOC_WORDS):
        return "toc", 0.9, meta
    if idx == 0:
        return "cover", 0.9, meta
    if has_chart:
        return "bar_chart", 0.85, meta
    if has_table:
        return "table", 0.85, meta
    # 章节过渡页：大数字 + 短标题
    if any(re.fullmatch(r"0?\d{1,2}", t) for t in texts) and body_chars <= 30:
        return "section", 0.8, meta
    # 内容框架页：有标题区、正文区干净（有效文字很少）——克隆底版首选
    if meta["title_geo"] is not None and body_chars <= 30:
        return "content_frame", 0.85, meta
    if body_chars <= 15 and len(shapes) <= 4:
        return "content_frame", 0.6, meta
    # 卡片页：3~4 个尺寸相近的文本框
    tboxes = [s for s in shapes if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()]
    if 3 <= len(tboxes) - 1 <= 4:
        widths = sorted(s.width for s in tboxes if s.width)
        if widths and widths[-1] - widths[0] < widths[0] * 0.3:
            return ("three_cards" if len(tboxes) - 1 == 3 else "four_cards"), 0.6, meta
    if pictures and tboxes:
        return "image_text", 0.55, meta
    if tboxes:
        return "title_content", 0.5, meta
    return "unknown", 0.3, meta
