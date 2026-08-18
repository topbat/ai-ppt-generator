from pptx import Presentation

from app.ppt.layout_recipes import RECIPE_CATALOG, project_recipe
from app.ppt.recipe_renderer import paint_recipe
from app.ppt.renderer import RenderContext, render_presentation
from app.schemas.composition import Box, TemplateSpaceContract
from app.schemas.presentation import PresentationSpec, SlideSpec, ThemeSpec


def _composition():
    contract = TemplateSpaceContract(
        page_width=13.33,
        page_height=7.5,
        safe_zone=Box(x=0.7, y=1.4, width=11.93, height=5.5),
    )
    recipe = next(recipe for recipe in RECIPE_CATALOG if recipe.id == "focus_statement")
    projected = project_recipe(recipe, contract)
    return contract, projected


def test_recipe_renderer_keeps_generated_shapes_inside_safe_zone():
    contract, recipe = _composition()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rc = RenderContext(ThemeSpec(), "medium")
    data = {
        "page": 1,
        "type": "title_content",
        "elements": [{"type": "bullet_group", "items": ["核心结论", "支持事实"]}],
        "visual_plan": {"regions": [region.model_dump() for region in recipe.regions]},
        "layout_recipe": recipe.id,
    }

    assert paint_recipe(rc, slide, data) is True
    safe = contract.safe_zone
    for shape in slide.shapes:
        assert shape.left.inches >= safe.x - 0.01
        assert shape.top.inches >= safe.y - 0.01
        assert shape.left.inches + shape.width.inches <= safe.x + safe.width + 0.01
        assert shape.top.inches + shape.height.inches <= safe.y + safe.height + 0.01


def test_rendered_pptx_contains_structured_speaker_notes(tmp_path):
    _, recipe = _composition()
    spec = PresentationSpec(
        title="备注测试",
        total_pages=1,
        slides=[
            SlideSpec(
                page=1,
                type="title_content",
                title="结论",
                elements=[{"type": "bullet_group", "items": ["主要观点"]}],
                layout_recipe=recipe.id,
                visual_plan={"regions": [region.model_dump() for region in recipe.regions]},
                speaker_notes={"details": ["补充材料一", "补充材料二"], "sources": []},
            )
        ],
    )
    path = tmp_path / "notes.pptx"

    render_presentation(spec, str(path))
    reopened = Presentation(path)

    assert "补充材料一" in reopened.slides[0].notes_slide.notes_text_frame.text
