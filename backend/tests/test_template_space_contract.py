from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.parser.template_parser import _space_contract
from app.schemas.composition import TemplateSpaceContract


def test_space_contract_reserves_title_logo_and_footer(tmp_path):
    logo_path = tmp_path / "logo.png"
    Image.new("RGB", (80, 30), "#1677ff").save(logo_path)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.6)).text = "页面标题"
    slide.shapes.add_picture(str(logo_path), Inches(11.7), Inches(0.35), Inches(0.8), Inches(0.3))
    slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(5), Inches(0.2)).text = "公司机密"

    contract = _space_contract(
        slide,
        prs,
        {"left": 731520, "top": 365760, "width": 7315200, "height": 548640},
    )
    validated = TemplateSpaceContract.model_validate(contract)

    assert validated.safe_zone.y >= 1.28
    assert validated.safe_zone.y + validated.safe_zone.height <= 7.0
    assert {zone.role for zone in validated.protected_zones} == {"logo", "footer"}


def test_space_contract_is_attached_to_every_parsed_slide(tmp_path):
    path = tmp_path / "template.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)

    from app.parser.template_parser import parse_template

    parsed = parse_template(str(path))

    assert parsed["slides"]
    assert all("space_contract" in slide["layout_meta"] for slide in parsed["slides"])


def test_space_contract_ignores_full_page_shape_misidentified_as_title():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    contract = _space_contract(
        slide,
        prs,
        {"left": 0, "top": 0, "width": int(prs.slide_width), "height": int(prs.slide_height)},
    )

    assert contract["safe_zone"]["y"] == 0.7
    assert contract["safe_zone"]["y"] + contract["safe_zone"]["height"] <= 7.5
