from collections import Counter
from types import SimpleNamespace

from pptx import Presentation

from app.pipeline.guards.capacity_guard import fit_page_capacity
from app.pipeline.stages.key_slide_stage import KeySlideDesignStage
from app.ppt.deck_rhythm import layout_fingerprint, rebalance_families, rhythm_issues
from app.ppt.layout_guard import guard_regions
from app.ppt.layout_recipes import RECIPE_CATALOG, project_recipe
from app.ppt.renderer import render_presentation
from app.schemas.composition import Box, SlideVisualPlan, TemplateSpaceContract
from app.schemas.presentation import PresentationSpec, SlideSpec


class RejectingGateway:
    def chat_json(self, *args, **kwargs):
        return {"invalid": True}


class NullStorage:
    def put_json(self, key, value):
        raise AssertionError("agent failure must preserve ordinary PresentationSpec")


def _contract():
    return TemplateSpaceContract(
        page_width=13.333,
        page_height=7.5,
        safe_zone=Box(x=0.7, y=1.4, width=11.933, height=5.5),
    )


def test_enterprise_composition_end_to_end_and_agent_failure_fallback(tmp_path):
    contract = _contract()
    families = ["editorial", "split", "cards_grid", "flow", "structure", "data_focus", "single_focus"]
    plans = [
        SlideVisualPlan(
            page=page,
            narrative_role="conclusion" if page in {4, 8, 12} else "evidence",
            importance="primary" if page % 3 == 0 else "secondary",
            thesis=f"第{page}页核心判断",
            layout_family=families[(page - 1) % len(families)],
            focal_element="bullet_group",
            focal_position=["left", "right", "center"][(page - 1) % 3],
            key_slide_candidate=page in {4, 8, 12},
        )
        for page in range(1, 13)
    ]
    candidates = {
        plan.page: [
            project_recipe(recipe, contract)
            for recipe in RECIPE_CATALOG
            if "bullet_group" in recipe.allowed_elements
            and guard_regions(project_recipe(recipe, contract), contract) == []
        ]
        for plan in plans
    }
    selected = rebalance_families(plans, candidates, "enterprise-e2e")

    assert rhythm_issues(plans, selected) == []
    counts = Counter(recipe.family for recipe in selected.values())
    assert max(counts.values()) / len(plans) <= 0.30
    fingerprints = [layout_fingerprint(plan, selected[plan.page]) for plan in plans]
    assert all(left != right for left, right in zip(fingerprints, fingerprints[1:]))

    slides = []
    preserved_numbers = set()
    for plan in plans:
        content = {
            "title": plan.thesis,
            "elements": [
                {
                    "type": "bullet_group",
                    "items": (
                        ["已覆盖23个系统", "主要结论", "支持材料一", "支持材料二"]
                        if plan.page == 1
                        else [f"第{plan.page}页主要观点", "支持依据"]
                    ),
                }
            ],
        }
        numbers = {"23"} if plan.page == 1 else set()
        fitted = fit_page_capacity(content, max_points=2, max_chars=80, source_numbers=numbers)
        preserved_numbers.update(fitted.fact_numbers)
        recipe = selected[plan.page]
        slides.append(
            SlideSpec(
                page=plan.page,
                chapter_id=f"ch{plan.page}",
                type="title_content",
                title=plan.thesis,
                elements=fitted.visible["elements"],
                layout_recipe=recipe.id,
                visual_plan={**plan.model_dump(), "regions": [region.model_dump() for region in recipe.regions]},
                speaker_notes=fitted.notes if fitted.notes["moved_count"] else None,
            )
        )
    spec = PresentationSpec(title="企业构图回归", total_pages=12, mode="premium", slides=slides)
    ordinary_before = spec.model_dump()

    ctx = SimpleNamespace(
        job_pk=77,
        biz_id="enterprise-e2e",
        storage=NullStorage(),
        _cache={"presentation_spec": spec},
        data={
            "LAYOUT": {"pjson_key": "unused.json"},
            "PARSE_TPL": {"space_contract": contract.model_dump(), "design_tokens": {}},
            "PLAN": {"plan": [{"page": page, "chapter_idx": page, "type": "title_content"} for page in range(1, 13)]},
            "STORYBOARD": {"slides": [plan.model_dump() for plan in plans]},
            "COMPOSE": {"composition_by_page": {str(page): {"family": selected[page].family} for page in selected}},
        },
    )
    key_result = KeySlideDesignStage(RejectingGateway()).run(ctx)

    assert key_result["selected"]
    assert key_result["applied"] == []
    assert key_result["fallback"] == key_result["selected"]
    assert spec.model_dump() == ordinary_before

    out_path = tmp_path / "enterprise-e2e.pptx"
    render_presentation(spec, str(out_path))
    reopened = Presentation(out_path)

    assert len(reopened.slides) == 12
    assert preserved_numbers == {"23"}
    first_notes = reopened.slides[0].notes_slide.notes_text_frame.text
    assert "23" in " ".join([slides[0].elements[0]["items"][0], first_notes])
    for slide in reopened.slides:
        for shape in slide.shapes:
            top = shape.top.inches
            bottom = top + shape.height.inches
            if top < contract.safe_zone.y - 0.01 or top > contract.safe_zone.y + contract.safe_zone.height:
                continue
            assert shape.left.inches >= contract.safe_zone.x - 0.02
            assert shape.left.inches + shape.width.inches <= contract.safe_zone.x + contract.safe_zone.width + 0.02
            assert bottom <= contract.safe_zone.y + contract.safe_zone.height + 0.02
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip() and run.font.size:
                            assert run.font.size.pt >= 16
