"""视觉优化冒烟测试（不依赖 DB/Redis/MinIO/LLM）：
1. 九维评分引擎结构与取值合法性；
2. Fix Ops 推导（低密度升格/两栏配平/超预算精简）与受控 DSL 校验；
3. 确定性调整器 + PPTX 几何微调；
4. 完整 BEAUTIFY 闭环（评分 → Ops → 调整 → 重渲染 → 复评，分数必须不降）。
运行：python tests/test_visual.py
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from app.ppt import design_tokens as dt  # noqa: E402
from app.ppt.renderer import render_presentation  # noqa: E402
from app.ppt.visual_ops import (FIX_OPS, SPEC_OPS, apply_spec_ops, derive_ops,  # noqa: E402
                                polish_pptx, validate_ops)
from app.ppt.visual_score import DIM_MAX, info_points, score_presentation  # noqa: E402
from app.schemas.presentation import PresentationSpec, SlideSpec, ThemeSpec  # noqa: E402

OUT_DIR = os.path.dirname(os.path.abspath(__file__))


def _bad_spec() -> PresentationSpec:
    """刻意制造视觉缺陷的演示文稿：低密度页 / 两栏失衡页 / 超预算拥挤页。"""
    slides = [
        SlideSpec(page=1, type="cover", title="视觉优化测试文稿", subtitle="BEAUTIFY 闭环"),
        # 低密度：只有 1 个信息点的 title_content → 应升格为大数字页（含数字）
        SlideSpec(page=2, type="title_content", title="市场规模",
                  key_message="市场空间巨大",
                  elements=[{"type": "bullet_group", "items": ["市场规模达到 1200亿元"]}]),
        # 低密度：无数字 → 应升格为大标题页（quote）
        SlideSpec(page=3, type="title_content", title="核心理念",
                  key_message="让数据驱动每一个决策",
                  elements=[{"type": "bullet_group", "items": ["以客户为中心"]}]),
        # 两栏失衡：前两条特别长、后一条极短
        SlideSpec(page=4, type="two_column", title="现状分析", elements=[
            {"type": "bullet_group", "items": [
                "当前各业务系统烟囱林立，数据孤岛问题严重，跨部门协同需要大量人工对账与重复录入",
                "核心流程仍依赖线下审批与纸质单据流转，平均处理时长超过五个工作日，客户体验持续下滑",
                "运维成本高"]}]),
        # 超预算拥挤：8 条长要点 → 应精简至预算恢复字号
        SlideSpec(page=5, type="title_content", title="重点工作", elements=[
            {"type": "bullet_group", "items": [
                f"第{i}项重点工作任务：包括方案设计、资源协调、里程碑管控与验收评估等多个环节的推进"
                for i in range(1, 9)]}]),
        SlideSpec(page=6, type="summary", title="总结", elements=[
            {"type": "bullet_group", "items": ["统一数字底座", "三年三步走", "价值可量化"]}]),
    ]
    return PresentationSpec(title="视觉优化测试", total_pages=len(slides),
                            theme=ThemeSpec(), slides=slides)


def test_tokens_grid():
    x, w = dt.col(0, 6)
    x2, w2 = dt.col(6, 6)
    assert abs((x2 + w2) - (dt.CANVAS_W - dt.MARGIN_X)) < 0.01, "12列网格右缘应对齐安全区"
    assert abs(x - dt.ANCHOR_X) < 0.001, "首列左缘应等于对齐锚线"
    assert abs(dt.snap8(1.0) * 72 % 8) < 0.001, "snap8 应吸附 8pt 倍数"
    print("✓ Design Token 12列网格/锚线/8pt吸附 通过")


def _check_score_engine():
    spec = _bad_spec()
    out = os.path.join(OUT_DIR, "visual_bad.pptx")
    result = render_presentation(spec, out)
    assert not result["degraded_pages"], f"渲染存在降级页: {result['degraded_pages']}"

    score = score_presentation(spec, out, result["notes"])
    assert sum(DIM_MAX.values()) == 100, "九维满分合计应为 100"
    assert 0 <= score["score"] <= 100
    assert len(score["pages"]) == len(spec.slides)
    for p in score["pages"]:
        for dim, v in p["dimensions"].items():
            assert 0 <= v <= DIM_MAX[dim], f"P{p['page']} 维度 {dim}={v} 越界"
    # 刻意制造的缺陷必须被识别：低密度（P2/P3）与超预算（P5）
    deds = {p["page"]: " ".join(d["detail"] for d in p["deductions"]) for p in score["pages"]}
    assert "低密度" in deds[2] and "低密度" in deds[3], f"低密度页未被识别: {deds}"
    assert "超出版式预算" in deds[5], f"拥挤页未被识别: {deds[5]}"
    assert info_points(spec.slides[4].model_dump()) == 8
    print(f"✓ 九维评分引擎通过（整册 {score['score']} 分，"
          f"P2={deds[2][:18]}… P5={deds[5][:18]}…）")
    return spec, out, result, score


def _check_ops_dsl(spec, score):
    ops = derive_ops(score, spec.model_dump())
    kinds = {o["op"] for o in ops}
    assert kinds <= set(FIX_OPS), f"产出了枚举外 op: {kinds - set(FIX_OPS)}"
    by_page = {(o["op"], o["page"]) for o in ops}
    assert ("upgrade_layout", 2) in by_page, f"P2 未推导出升格指令: {ops}"
    assert ("upgrade_layout", 3) in by_page, f"P3 未推导出升格指令: {ops}"
    assert ("rebalance", 4) in by_page, f"P4 未推导出两栏配平指令: {ops}"
    assert ("font_token", 5) in by_page, f"P5 未推导出精简指令: {ops}"

    # 受控 DSL：枚举外/非法页码一律拒绝
    valid = validate_ops([{"op": "delete_page", "page": 2},
                          {"op": "upgrade_layout", "page": 99},
                          {"op": "rebalance", "page": 4, "detail": "ok"}], total_pages=6)
    assert len(valid) == 1 and valid[0]["op"] == "rebalance", f"DSL 校验失败: {valid}"
    print(f"✓ Fix Ops DSL 通过（推导 {len(ops)} 条指令，枚举外/非法页码已拒绝）")
    return ops


def _check_adjuster_and_loop(spec, pptx0, score0, ops):
    spec_ops = [o for o in ops if o["op"] in SPEC_OPS]
    new_dict, applied, skipped = apply_spec_ops(spec.model_dump(), spec_ops, "medium")
    assert applied, "确定性调整器未应用任何指令"
    pages = {s["page"]: s for s in new_dict["slides"]}
    assert pages[2]["type"] == "key_number", f"P2 应升格为大数字页，实际 {pages[2]['type']}"
    assert pages[3]["type"] == "quote", f"P3 应升格为大标题页，实际 {pages[3]['type']}"
    assert len(pages[4]["elements"]) == 2, "P4 两栏应配平为两个要点组"
    p5_items = pages[5]["elements"][0]["items"]
    assert len(p5_items) <= 5 and all(len(i) <= 40 for i in p5_items), "P5 应精简至密度预算"

    # 重渲染 → 几何微调 → 复评（闭环收益必须为正）
    new_spec = PresentationSpec.model_validate(new_dict)
    out1 = os.path.join(OUT_DIR, "visual_fixed.pptx")
    r1 = render_presentation(new_spec, out1)
    stats = polish_pptx(out1)
    score1 = score_presentation(new_spec, out1, r1["notes"])
    assert score1["score"] > score0["score"], \
        f"闭环后分数未提升: {score0['score']} → {score1['score']}"
    from pptx import Presentation
    assert len(Presentation(out1).slides._sldIdLst) == len(new_spec.slides), "微调后文件应完整"
    print(f"✓ BEAUTIFY 闭环通过（视觉分 {score0['score']} → {score1['score']}，"
          f"应用 {len(applied)} 条 Ops，几何微调 对齐{stats['aligned']}/吸附{stats['spacing_snapped']}）")


def test_good_deck_stability():
    """常规完整版式（沿用 test_smoke 结构）分数应处于合理高位，且推导不出升格类指令。"""
    slides = [
        SlideSpec(page=1, type="cover", title="企业数字化转型建设方案", subtitle="2026年度汇报"),
        SlideSpec(page=2, type="three_cards", title="三大目标", elements=[
            {"type": "cards", "items": [{"title": "降本", "desc": "统一平台减少重复建设"},
                                        {"title": "增效", "desc": "流程自动化提升效率"},
                                        {"title": "创新", "desc": "数据驱动业务创新"}]}]),
        SlideSpec(page=3, type="key_number", title="核心指标", elements=[
            {"type": "key_number", "items": [{"value": "23个", "label": "业务系统"},
                                             {"value": "1.2亿", "label": "总投资"},
                                             {"value": "3年", "label": "建设周期"}]}]),
        SlideSpec(page=4, type="title_content", title="行业趋势", elements=[
            {"type": "bullet_group", "items": ["政策驱动数字化提速", "行业竞争加剧",
                                               "客户体验要求提升", "技术栈快速演进"]}]),
    ]
    spec = PresentationSpec(title="常规文稿", total_pages=len(slides),
                            theme=ThemeSpec(), slides=slides)
    out = os.path.join(OUT_DIR, "visual_good.pptx")
    r = render_presentation(spec, out)
    polish_pptx(out)
    score = score_presentation(spec, out, r["notes"])
    assert score["score"] >= 80, f"常规文稿视觉分异常偏低: {score['score']}"
    ops = derive_ops(score, spec.model_dump())
    assert not any(o["op"] == "upgrade_layout" for o in ops), f"常规文稿不应触发升格: {ops}"
    print(f"✓ 常规文稿稳定性通过（视觉分 {score['score']}，无升格类误报）")


def test_score_engine():
    _check_score_engine()


def test_ops_dsl():
    spec, _, _, score = _check_score_engine()
    _check_ops_dsl(spec, score)


def test_adjuster_and_loop():
    spec, pptx0, _, score0 = _check_score_engine()
    ops = _check_ops_dsl(spec, score0)
    _check_adjuster_and_loop(spec, pptx0, score0, ops)


if __name__ == "__main__":
    test_tokens_grid()
    spec, pptx0, render_result, score0 = _check_score_engine()
    ops = _check_ops_dsl(spec, score0)
    _check_adjuster_and_loop(spec, pptx0, score0, ops)
    test_good_deck_stability()
    print("\n视觉优化全部冒烟测试通过 ✅")
